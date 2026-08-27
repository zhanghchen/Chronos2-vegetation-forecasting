# -*- coding: utf-8 -*-
"""Short teaching deck accompanying TUTORIAL.md - for walking another
researcher through Chronos-2 in a meeting rather than having them read the
full markdown cold. Reuses this project's established slide-deck visual
identity."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

ROOT = str(Path(__file__).resolve().parent)

INK = RGBColor(0x1C, 0x21, 0x19)
MUTED = RGBColor(0x5C, 0x63, 0x55)
FAINT = RGBColor(0x8B, 0x93, 0x82)
ACCENT = RGBColor(0x2F, 0x6F, 0x5E)
ACCENT_DARK = RGBColor(0x1E, 0x4A, 0x33)
ACCENT_TINT = RGBColor(0xE7, 0xF0, 0xE6)
WARN = RGBColor(0xB5, 0x65, 0x1D)
WARN_TINT = RGBColor(0xFB, 0xEE, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDE, 0xE2, 0xD6)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.25)
TAKEAWAY_H = Inches(0.78)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = RULE; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.12
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic
        r.font.name = FONT
    return box


def add_bullets(slide, left, top, width, height, items, size=17, color=INK, space_after=10, bullet_color=ACCENT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.line_spacing = 1.08
        r0 = p.add_run(); r0.text = "▪  "; r0.font.size = Pt(size); r0.font.color.rgb = bullet_color; r0.font.name = FONT
        r1 = p.add_run(); r1.text = item; r1.font.size = Pt(size); r1.font.color.rgb = color; r1.font.name = FONT
    return box


def add_eyebrow(slide, text):
    add_text(slide, MARGIN, Inches(0.42), Inches(12), Inches(0.32), [(text, 13, ACCENT, True, False)])


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    add_text(slide, MARGIN, Inches(0.72), Inches(12.2), Inches(0.62), [(title, 24, INK, True, False)])
    add_rect(slide, MARGIN, Inches(1.38), Inches(1.0), Pt(3), ACCENT)


def add_takeaway(slide, text):
    add_rect(slide, 0, TAKEAWAY_TOP, SLIDE_W, TAKEAWAY_H, ACCENT_TINT)
    add_rect(slide, 0, TAKEAWAY_TOP, Inches(0.12), TAKEAWAY_H, ACCENT)
    add_text(slide, Inches(0.45), TAKEAWAY_TOP, SLIDE_W - Inches(0.9), TAKEAWAY_H,
              [("TAKEAWAY   ", 12, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.55), TAKEAWAY_TOP, SLIDE_W - Inches(2.0), TAKEAWAY_H,
              [(text, 15, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)


def page_num(slide, n):
    add_text(slide, SLIDE_W - Inches(0.7), Inches(0.42), Inches(0.4), Inches(0.3),
              [(str(n), 11, FAINT, False, False)], align=PP_ALIGN.RIGHT)


def content_box(has_takeaway=True):
    bottom = (TAKEAWAY_TOP - Inches(0.15)) if has_takeaway else (SLIDE_H - Inches(0.3))
    return CONTENT_TOP, bottom - CONTENT_TOP


def code_box(slide, left, top, width, height, code, size=13):
    shp = add_rect(slide, left, top, width, height, RGBColor(0xF5, 0xF6, 0xF1), line=True)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12); tf.margin_right = Pt(12); tf.margin_top = Pt(10); tf.margin_bottom = Pt(10)
    for i, line in enumerate(code.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.name = "Consolas"; r.font.color.rgb = INK
    return shp


def flow_box(slide, left, top, width, height, text, fill, text_color=WHITE, size=11):
    shp = add_rounded(slide, left, top, width, height, fill)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.color.rgb = text_color; r.font.name = FONT
    return shp


# ============================================================ SLIDE 1: TITLE
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), ACCENT)
add_text(s, MARGIN, Inches(0.28), Inches(12.2), Inches(0.6), [("A Practical Guide to Chronos-2", 26, WHITE, True, False)])
page_num(s, 1)
top, h = content_box()
add_text(s, MARGIN, top + Inches(0.2), Inches(11.5), Inches(0.6),
          [("Applying a pretrained time-series foundation model to a new dataset", 18, MUTED, True, True)])
add_bullets(s, MARGIN, top + Inches(1.1), Inches(11.5), Inches(2.5), [
    "Zero-shot forecasting, no training required to get a first baseline",
    "Past and known-future covariates (e.g. climate, calendar features)",
    "Optional fine-tuning: full-weight or LoRA",
    "A general, reusable data-preparation pattern that works across domains",
], size=18, space_after=14)
add_takeaway(s, "Full written tutorial: TUTORIAL.md — this deck is a walkthrough companion, not a replacement.")

# ============================================================ SLIDE 2: WHAT IS CHRONOS-2
s = add_slide(); set_bg(s)
add_title(s, "What Chronos-2 Does", eyebrow="BACKGROUND")
page_num(s, 2)
top, h = content_box()
header_items = [
    ("Zero-shot forecasting", "Forecast a new series with no training — load the pretrained weights and predict."),
    ("Covariates", "Add auxiliary series: past-only, or known-future (e.g. a weather forecast)."),
    ("Fine-tuning", "Optionally adapt the weights: full-weight or a small LoRA adapter."),
    ("Cross-learning", "Optionally let multiple series in a batch share information at inference/training time."),
]
y = top
for name, desc in header_items:
    add_rounded(s, MARGIN, y, Inches(2.6), Inches(0.9), ACCENT_TINT)
    add_text(s, MARGIN + Inches(0.15), y + Inches(0.12), Inches(2.3), Inches(0.65),
              [(name, 13.5, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MARGIN + Inches(2.85), y + Inches(0.05), Inches(9), Inches(0.85),
              [(desc, 15, INK, False, False)], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.05)
add_takeaway(s, "This tutorial focuses on the practical workflow, not the model's internal architecture.")

# ============================================================ SLIDE 3: INSTALL + LOAD
s = add_slide(); set_bg(s)
add_title(s, "Installation and Loading the Model", eyebrow="SETUP")
page_num(s, 3)
top, h = content_box()
code_box(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(1.5), """
pip install chronos-forecasting torch pandas numpy matplotlib scikit-learn
""".strip("\n"), size=15)
code_box(s, MARGIN, top + Inches(1.75), SLIDE_W - 2 * MARGIN, Inches(2.3), """
from chronos import BaseChronosPipeline
pipeline = BaseChronosPipeline.from_pretrained("amazon/chronos-2", device_map="cuda")

assert hasattr(pipeline, "predict_quantiles")
print(pipeline.model_context_length, pipeline.model_prediction_length)
""".strip("\n"), size=15)
add_takeaway(s, "BaseChronosPipeline.from_pretrained auto-dispatches to the right pipeline class — works on CPU or GPU.")

# ============================================================ SLIDE 4: DATA FORMAT
s = add_slide(); set_bg(s)
add_title(s, "Data Format: Four Roles", eyebrow="DATA PREPARATION")
page_num(s, 4)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(12), Inches(1.0), [
    "Every column in your dataset is one of four roles:",
], size=16, space_after=6)
header = ["timestamp", "series ID", "target", "covariates"]
colors = [MUTED, MUTED, ACCENT, ACCENT_DARK]
x = MARGIN
box_w = Inches(2.75)
for label, color in zip(header, colors):
    flow_box(s, x, top + Inches(0.7), box_w, Inches(0.8), label, color, size=15)
    x += box_w + Inches(0.15)

add_bullets(s, MARGIN, top + Inches(1.9), Inches(12), Inches(2.2), [
    "past_covariates: observed up through \"now\" — required for every covariate the model uses.",
    "future_covariates: values already known for the forecast horizon (weather forecast, calendar) — "
    "must be a subset of past_covariates' keys.",
    "Never put the target itself, or an unknown future value, into future_covariates.",
], size=16, space_after=10)
add_takeaway(s, "past_covariates vs. future_covariates is the single most important distinction to get right.")

# ============================================================ SLIDE 5: WORKFLOW
s = add_slide(); set_bg(s)
add_title(s, "Recommended Workflow", eyebrow="METHODOLOGY")
page_num(s, 5)
top, h = content_box()

steps = ["Inspect\ndataset", "Identify\ntarget &\ncovariates", "Align\ntemporal\nresolution",
         "Check\nmissing\nvalues", "Train/val/\ntest split", "Zero-shot\nChronos-2",
         "Evaluate", "Add useful\ncovariates", "Compare vs.\nzero-shot", "Fine-tune\n(if needed)",
         "Validate\ncarefully", "Test once,\nfinal"]
n = len(steps)
cols = 6
box_w = Inches(1.9)
box_h = Inches(0.85)
gap_x = (SLIDE_W - 2 * MARGIN - cols * box_w) // (cols - 1)
gap_y = Inches(0.35)
for i, step in enumerate(steps):
    row, col = divmod(i, cols)
    x = MARGIN + col * (box_w + gap_x)
    y = top + Inches(0.1) + row * (box_h + gap_y)
    color = ACCENT if row == 0 else ACCENT_DARK
    flow_box(s, x, y, box_w, box_h, f"{i+1}. {step}", color, size=11)
add_takeaway(s, "Only fine-tune after zero-shot has been evaluated and covariates have earned their place.")

# ============================================================ SLIDE 6: FINE-TUNING
s = add_slide(); set_bg(s)
add_title(s, "Fine-Tuning: pipeline.fit(...)", eyebrow="ADAPTATION")
page_num(s, 6)
top, h = content_box()
code_box(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(2.4), """
finetuned = pipeline.fit(
    inputs=[train_item],
    prediction_length=prediction_length,
    validation_inputs=[val_item],      # enables checkpoint selection
    finetune_mode="lora",              # "full" (default) or "lora"
    learning_rate=1e-5,
    num_steps=1000,
)
""".strip("\n"), size=14.5)
add_bullets(s, MARGIN, top + Inches(2.7), Inches(12), Inches(1.6), [
    "\"full\" updates every weight; \"lora\" trains a small adapter (needs peft) — both are officially supported.",
    "validation_inputs auto-enables load_best_model_at_end=True (best validation loss, not the final step).",
    "Trainable parameter count does not reliably predict fine-tuning quality — validate, don't assume.",
], size=15, space_after=8)
add_takeaway(s, "Always compare the fine-tuned result against the zero-shot baseline on the SAME validation window.")

# ============================================================ SLIDE 7: TRAIN/VAL/TEST
s = add_slide(); set_bg(s)
add_title(s, "Train / Validation / Test — Chronologically", eyebrow="EVALUATION PROTOCOL")
page_num(s, 7)
top, h = content_box()
add_rect(s, MARGIN, top + Inches(0.1), Inches(11.5), Inches(0.65), WARN_TINT)
add_text(s, MARGIN + Inches(0.25), top + Inches(0.18), Inches(11), Inches(0.5),
          [("Never split individual timestamps randomly — this leaks future information into training.", 15, WARN, True, False)],
          anchor=MSO_ANCHOR.MIDDLE)

rows = [("2000–2017", "→", "2018"), ("2000–2018", "→", "2019"), ("2000–2019", "→", "2020"), ("2000–2020", "→", "2021")]
y = top + Inches(1.1)
for train, arrow, val in rows:
    flow_box(s, MARGIN, y, Inches(2.6), Inches(0.55), train, ACCENT_TINT, text_color=ACCENT_DARK, size=13)
    add_text(s, MARGIN + Inches(2.7), y, Inches(0.5), Inches(0.55), [("→", 16, MUTED, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    flow_box(s, MARGIN + Inches(3.3), y, Inches(1.6), Inches(0.55), val, ACCENT, size=13)
    y += Inches(0.68)

add_text(s, MARGIN + Inches(5.6), top + Inches(1.1), Inches(0.6), Inches(2.7), [("then", 15, MUTED, True, True)], anchor=MSO_ANCHOR.MIDDLE)
flow_box(s, MARGIN + Inches(6.4), top + Inches(1.6), Inches(3.0), Inches(0.7), "Final refit: 2000–2021", ACCENT_DARK, size=13)
add_text(s, MARGIN + Inches(6.4), top + Inches(2.35), Inches(0.4), Inches(0.4), [("↓", 16, MUTED, True, False)])
flow_box(s, MARGIN + Inches(6.4), top + Inches(2.75), Inches(3.0), Inches(0.7), "Evaluate ONCE: 2022", WARN, size=13)

add_takeaway(s, "The final test year is touched exactly once — after every modeling decision is already locked in.")

# ============================================================ SLIDE 8: PITFALLS
s = add_slide(); set_bg(s)
add_title(s, "Common Pitfalls", eyebrow="TROUBLESHOOTING")
page_num(s, 8)
top, h = content_box()
add_bullets(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.6), [
    "Resolution mismatch: context_length / prediction_length are counts of TIME STEPS, not days or months.",
    "Off-by-one splits: a date-based cutoff can silently produce the wrong number of future rows — "
    "split by exact row count when precision matters.",
    "predict_df requires strictly regular timestamps; irregular data needs the lower-level dict API "
    "(or align_frequency() first).",
    "Static features (vegetation type, patient category, sensor model) broadcast as a constant covariate "
    "often get erased by per-series normalization — see next section of TUTORIAL.md.",
    "\"More covariates\" is not automatically better — validate each addition against the zero-shot "
    "baseline before trusting it.",
], size=16, space_after=11)
add_takeaway(s, "Full troubleshooting table (symptom → cause → check → fix) in TUTORIAL.md Part 15.")

# ============================================================ SLIDE 9: STATIC FEATURES
s = add_slide(); set_bg(s)
add_title(s, "Static Features Need Special Attention", eyebrow="MODELING LESSON")
page_num(s, 9)
top, h = content_box()

box_w2, box_h2 = Inches(2.9), Inches(0.8)
gap2 = Inches(0.5)
x0 = MARGIN + Inches(0.4)
flow_box(s, x0, top + Inches(0.3), box_w2, box_h2, "Static value\n[0.65, 0.65, 0.65, ...]", MUTED, size=13)
flow_box(s, x0 + box_w2 + gap2, top + Inches(0.3), box_w2, box_h2, "Per-series\nnormalization", WARN, size=13)
flow_box(s, x0 + 2 * (box_w2 + gap2), top + Inches(0.3), box_w2, box_h2, "Zero variance\n→ signal erased", ACCENT_DARK, size=13)
for i in range(2):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        x0 + (i + 1) * box_w2 + i * gap2, top + Inches(0.7), x0 + (i + 1) * box_w2 + (i + 1) * gap2, top + Inches(0.7))
    conn.line.color.rgb = FAINT; conn.line.width = Pt(2.5)

add_bullets(s, MARGIN, top + Inches(1.55), SLIDE_W - 2 * MARGIN, Inches(2.6), [
    "A feature that describes a whole series (vegetation type, patient group, sensor model) is NOT an "
    "ordinary time-varying covariate — check whether your model's normalization operates per-series "
    "before assuming a naively broadcast static feature will be used.",
    "Chronos-2 supports categorical covariates that CHANGE over time (day-of-week, a shifting regime) — "
    "this is a different feature and does not solve the static-feature problem.",
    "Verify directly: perturb the static value while holding everything else fixed, and check whether "
    "the forecast actually changes.",
], size=15.5, space_after=10)
add_takeaway(s, "Inspect how static covariates are represented and normalized before assuming the model can use them.")

# ============================================================ SLIDE 10: WHERE TO GO NEXT
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.5), ACCENT)
add_text(s, MARGIN, Inches(0.35), Inches(12.2), Inches(0.9), [("Getting Started", 24, WHITE, True, False)])
page_num(s, 10)
top = Inches(1.85)
add_bullets(s, MARGIN, top, Inches(12.2), Inches(3.2), [
    "TUTORIAL.md — the full 17-part written tutorial, with a verified API reference at the end.",
    "chronos2_tutorial.ipynb — run it end-to-end on the included synthetic dataset.",
    "chronos2_template.py — copy it, edit the 6-line CONFIG block, and point it at your own data.",
    "example_config.py — a worked example adapting the template to a different domain.",
], size=17, space_after=14, bullet_color=ACCENT)
add_takeaway(s, "Start with zero-shot on your own data before adding covariates or fine-tuning.")

out_path = f"{ROOT}/Chronos2_Tutorial_Deck.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
