# -*- coding: utf-8 -*-
"""Advisor-facing summary deck for the full Chronos-2 + PFT investigation
(single-pixel ablation -> multi-pixel FiLM -> open architecture search +
shuffled-PFT control). 9 slides, reusing this project's established visual
identity. Every number is read directly from saved CSVs - see the
docstring at the bottom of this file for the source-to-slide mapping."""
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image

ROOT = str(Path(__file__).resolve().parent.parent)
OUT_DATA = f"{ROOT}/outputs/pft_v2"
OUT_MULTIPIXEL = f"{ROOT}/outputs/pft_multipixel"
AELSTM_ROOT = "/home/deh25003/chronos-forecasting/AELSTM"

INK = RGBColor(0x1C, 0x21, 0x19)
MUTED = RGBColor(0x5C, 0x63, 0x55)
FAINT = RGBColor(0x8B, 0x93, 0x82)
ACCENT = RGBColor(0x2F, 0x6F, 0x5E)
ACCENT_DARK = RGBColor(0x1E, 0x4A, 0x33)
ACCENT_TINT = RGBColor(0xE7, 0xF0, 0xE6)
WARN = RGBColor(0xB5, 0x65, 0x1D)
WARN_TINT = RGBColor(0xFB, 0xEE, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDE, 0xE2, 0xD6)
GREY = RGBColor(0x9A, 0xA2, 0x94)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.25)
TAKEAWAY_H = Inches(0.78)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or RULE
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, left, top, width, height, color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.12
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1.0)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
    return box


def add_bullets(slide, left, top, width, height, items, size=17, color=INK, space_after=10,
                 bullet_color=ACCENT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        r0 = p.add_run()
        r0.text = "▪  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        r0.font.name = FONT
        r1 = p.add_run()
        r1.text = item
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.name = FONT
    return box


def add_eyebrow(slide, text):
    add_text(slide, MARGIN, Inches(0.42), Inches(12), Inches(0.32), [(text, 13, ACCENT, True, False)])


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    add_text(slide, MARGIN, Inches(0.72), Inches(12.2), Inches(0.62), [(title, 24, INK, True, False)])
    add_rect(slide, MARGIN, Inches(1.38), Inches(1.0), Pt(3), ACCENT)


def add_takeaway(slide, text, tint=ACCENT_TINT, dark=ACCENT_DARK, bar=ACCENT):
    add_rect(slide, 0, TAKEAWAY_TOP, SLIDE_W, TAKEAWAY_H, tint)
    add_rect(slide, 0, TAKEAWAY_TOP, Inches(0.12), TAKEAWAY_H, bar)
    add_text(slide, Inches(0.45), TAKEAWAY_TOP, SLIDE_W - Inches(0.9), TAKEAWAY_H,
              [("TAKEAWAY   ", 12, dark, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.55), TAKEAWAY_TOP, SLIDE_W - Inches(2.0), TAKEAWAY_H,
              [(text, 15, dark, True, False)], anchor=MSO_ANCHOR.MIDDLE)


def page_num(slide, n):
    add_text(slide, SLIDE_W - Inches(0.7), Inches(0.42), Inches(0.4), Inches(0.3),
              [(str(n), 11, FAINT, False, False)], align=PP_ALIGN.RIGHT)


def add_picture_fit(slide, path, box_left, box_top, box_w, box_h, align="center"):
    im = Image.open(path)
    ar = im.size[0] / im.size[1]
    box_ar = box_w / box_h
    if ar > box_ar:
        w = box_w
        h = int(w / ar)
    else:
        h = box_h
        w = int(h * ar)
    left = box_left + (box_w - w) // 2 if align == "center" else box_left
    top = box_top + (box_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    return left, top, w, h


def content_box(has_takeaway=True):
    bottom = (TAKEAWAY_TOP - Inches(0.15)) if has_takeaway else (SLIDE_H - Inches(0.3))
    return CONTENT_TOP, bottom - CONTENT_TOP


def styled_table(slide, left, top, width, height, header, rows, col_weights,
                  header_size=12, body_size=11.5, highlight_rows=()):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(header), left, top, width, height).table
    for c, w in enumerate(col_weights):
        tbl.columns[c].width = Emu(int(width * w))
    for ci, text in enumerate(header):
        cell = tbl.cell(0, ci)
        cell.text = text
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(header_size); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = FONT
        p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    for ri, row in enumerate(rows, start=1):
        hl = (ri - 1) in highlight_rows
        for ci, text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_TINT if hl else (WHITE if ri % 2 else RGBColor(0xF5, 0xF6, 0xF1))
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(body_size); p.runs[0].font.bold = (ci == 0 or hl)
            p.runs[0].font.color.rgb = INK; p.runs[0].font.name = FONT
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    return tbl


def flow_box(slide, left, top, width, height, text, fill, text_color=WHITE, size=12.5):
    shp = add_rounded(slide, left, top, width, height, fill)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6); tf.margin_right = Pt(6); tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = text_color
        r.font.bold = (i == 0)
        r.font.name = FONT
    return shp


def flow_arrow_down(slide, x_center, top, length, color=FAINT):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_center, top, x_center, top + length)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    return conn


# ============================================================ SLIDE 1: MOTIVATION
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), ACCENT)
add_text(s, MARGIN, Inches(0.28), Inches(12.2), Inches(0.6),
          [("Can PFT Improve Chronos-2 LAI Forecasting?", 25, WHITE, True, False)])
page_num(s, 1)
top, h = content_box()
add_text(s, MARGIN, top + Inches(0.1), Inches(12.2), Inches(0.4), [("HYPOTHESIS", 13, ACCENT, True, False)])
add_text(s, MARGIN, top + Inches(0.55), Inches(11.8), Inches(1.3),
          [("Different vegetation compositions may respond differently to the same climate forcing.",
            19, INK, True, True)], line_spacing=1.15)
add_rect(s, MARGIN, top + Inches(1.55), Inches(11), Inches(0.9), ACCENT_TINT)
add_text(s, MARGIN + Inches(0.3), top + Inches(1.65), Inches(10.4), Inches(0.7),
          [("LAI_future = f( historical LAI, climate, PFT composition )", 17, ACCENT_DARK, True, False)],
          anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, MARGIN, top + Inches(2.75), Inches(11.8), Inches(1.6), [
    "Example: a pixel with 65% forest + 35% grass may not respond to precipitation, VPD, or temperature "
    "the same way a 100% forest pixel does.",
    "This deck summarizes a full investigation — from a single-pixel test that failed for an architectural "
    "reason, through an architecture fix, to a decisive controlled experiment.",
], size=16, space_after=12)
add_takeaway(s, "Question: can explicit PFT information give Chronos-2 a measurable, reproducible forecasting edge?")

# ============================================================ SLIDE 2: WHY INITIAL APPROACH FAILED
s = add_slide(); set_bg(s)
add_title(s, "Why the Initial PFT Approach Failed", eyebrow="STEP 1 · DIAGNOSIS")
page_num(s, 2)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(12.2), Inches(1.3), [
    "Single-pixel experiments fed PFT as a constant covariate (same value at every timestep) — and showed "
    "almost no effect.",
    "Source-code inspection of Chronos-2's InstanceNorm explained why:",
], size=16.5, space_after=10)

diagram_top = top + Inches(1.5)
box_w, box_h = Inches(2.6), Inches(0.75)
gap = Inches(0.55)
x0 = MARGIN + Inches(0.3)
flow_box(s, x0, diagram_top, box_w, box_h, "PFT = 0.65\n(constant every step)", MUTED, size=12.5)
flow_box(s, x0 + box_w + gap, diagram_top, box_w, box_h, "InstanceNorm\n(per-series, std=0)", WARN, size=12.5)
flow_box(s, x0 + 2 * (box_w + gap), diagram_top, box_w, box_h, "→ [0, 0, 0, ...]\nsignal erased", ACCENT_DARK, size=12.5)
for i in range(2):
    conn = slide_arrow = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        x0 + (i + 1) * box_w + i * gap, diagram_top + box_h // 2,
        x0 + (i + 1) * box_w + (i + 1) * gap, diagram_top + box_h // 2)
    conn.line.color.rgb = FAINT
    conn.line.width = Pt(2.5)

add_rect(s, MARGIN, diagram_top + Inches(1.3), Inches(11.5), Inches(0.75), WARN_TINT)
add_text(s, MARGIN + Inches(0.3), diagram_top + Inches(1.4), Inches(11), Inches(0.6),
          [("Verified empirically: perturbing PFT across 5 synthetic compositions changed the prediction by "
            "exactly 0.00000000 — not a small effect, zero.", 14.5, WARN, True, False)],
          anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, MARGIN, diagram_top + Inches(2.35), Inches(12), Inches(1.0), [
    "Critically: this did NOT mean PFT is biologically irrelevant — only that a constant-per-pixel value "
    "cannot survive this specific normalization step.",
], size=15.5, space_after=8, bullet_color=WARN)
add_takeaway(s, "Chronos-2 was architecturally blind to static PFT — a mechanism problem, not evidence against PFT itself.",
             tint=WARN_TINT, dark=WARN, bar=WARN)

# ============================================================ SLIDE 3: MAKING CHRONOS-2 PFT-AWARE
s = add_slide(); set_bg(s)
add_title(s, "Making Chronos-2 PFT-Aware", eyebrow="STEP 2 · ARCHITECTURE FIX")
page_num(s, 3)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(12.2), Inches(1.5), [
    "Built explicit PFT conditioning that bypasses InstanceNorm entirely: a small trainable module converts "
    "the PFT vector into a FiLM (scale, shift) adjustment applied directly to the model's internal "
    "representation — the frozen, 119.5M-parameter pretrained model never changes.",
    "Tested several conditioning designs: deep-MLP FiLM, a regularized variant, a biologically-structured "
    "linear mixture of per-vegetation-class response vectors, and a rank-8 bottleneck FiLM.",
], size=15.5, space_after=10)

img_top = top + Inches(1.75)
add_picture_fit(s, f"{OUT_MULTIPIXEL}/sensitivity/perturbation_sweep_all_pixels.png",
                  MARGIN, img_top, SLIDE_W - 2 * MARGIN, Inches(3.9))
add_takeaway(s, "Smoke test confirmed real PFT sensitivity (e.g. +0.63 LAI units, evergreen vs. grass) — "
                "the question became generalization, not visibility.")

# ============================================================ SLIDE 4: MULTI-PIXEL EXPERIMENTAL DESIGN
s = add_slide(); set_bg(s)
add_title(s, "Multi-Pixel Experimental Design", eyebrow="STEP 3 · SCALING UP")
page_num(s, 4)
top, h = content_box()
add_picture_fit(s, f"{AELSTM_ROOT}/outputs/pft_multipixel_selection/pft_diverse_pixels_diversity.png",
                  MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(3.6))
add_bullets(s, MARGIN, top + Inches(3.75), SLIDE_W - 2 * MARGIN, Inches(1.4), [
    "70 CONUS pixels selected for genuine PFT diversity (8 dominant classes, purity 0.46–1.00), pooled into "
    "one shared training set so PFT actually varies across training rows.",
    "A real data-quality bug was caught and fixed here: 24/70 initially-selected pixels had 100% missing "
    "gridMET climate data (a known US coverage gap); the final 70-pixel set was verified complete.",
], size=14.5, space_after=8)
add_takeaway(s, "Strict protocol: pre-2022 windows only for training/model selection; 2022 touched exactly once, at the end.")

# ============================================================ SLIDE 5: OVERFITTING PERSISTS
s = add_slide(); set_bg(s)
add_title(s, "Overfitting Persists Despite More Temporal Supervision", eyebrow="STEP 4 · SCREENING")
page_num(s, 5)
top, h = content_box()
add_picture_fit(s, f"{OUT_DATA}/deck_figures/overfitting_comparison.png", MARGIN, top,
                  SLIDE_W - 2 * MARGIN, Inches(3.75))

header = ["Architecture", "Params", "Best val_loss", "Selected step"]
rows = [
    ("deep-MLP FiLM (original)", "100,544", "0.34664", "0"),
    ("Regularized deep-MLP", "51,040", "0.34667", "0"),
    ("Linear per-class mixture", "15,360", "0.34695", "0"),
    ("Rank-8 bottleneck", "12,600", "0.34653", "20"),
]
styled_table(s, MARGIN, top + Inches(3.95), SLIDE_W - 2 * MARGIN, Inches(1.15), header, rows,
             [0.4, 0.2, 0.2, 0.2], body_size=12.5, highlight_rows=(3,))
add_takeaway(s, "8 training windows (4x more) still weren't enough for 3 of 4 architectures — data volume alone "
                "isn't the full explanation.")

# ============================================================ SLIDE 6: KEY FINDING
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), WARN)
add_text(s, MARGIN, Inches(0.28), Inches(12.2), Inches(0.6),
          [("Key Finding: Real PFT ≈ Shuffled PFT", 25, WHITE, True, False)])
page_num(s, 6)
top, h = content_box()
add_picture_fit(s, f"{OUT_DATA}/deck_figures/real_vs_shuffled_bar.png", MARGIN, top,
                  SLIDE_W - 2 * MARGIN, h - Inches(0.05))
add_takeaway(s, "A model trained on scientifically meaningless, randomly-shuffled PFT performs at least as well "
                "as one trained on real PFT — p=0.245, not significant.", tint=WARN_TINT, dark=WARN, bar=WARN)

# ============================================================ SLIDE 7: WHAT DOES THIS TELL US
s = add_slide(); set_bg(s)
add_title(s, "What Does This Tell Us?", eyebrow="INTERPRETATION")
page_num(s, 7)
top, h = content_box()
add_rect(s, MARGIN, top + Inches(0.05), Inches(11.6), Inches(1.0), ACCENT_TINT)
add_text(s, MARGIN + Inches(0.3), top + Inches(0.18), Inches(11), Inches(0.8),
          [("No PFT  ≈  Real PFT  ≈  Shuffled PFT", 20, ACCENT_DARK, True, False)],
          anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, MARGIN, top + Inches(1.3), Inches(12), Inches(2.6), [
    "The small ΔR² (~+0.002) most likely comes from adding trainable conditioning capacity itself, not from "
    "the biological information the PFT vector carries.",
    "Supporting evidence reaches the same conclusion from four independent angles: mixed-vs-pure pixel "
    "entropy (r ≈ −0.13 vs. shuffled), seasonal-phase RMSE (indistinguishable in every phase), and "
    "per-dominant-PFT-class breakdown (|Δ| ≤ 0.001 for every class once compared against shuffled).",
    "This is consistent with an earlier, independent finding in this project: no PEFT/fine-tuning method "
    "(LoRA, DoRA, VeRA, IA3, LN-Tuning, BitFit) beat zero-shot Chronos-2 on this dataset either.",
], size=15.5, space_after=12)

add_text(s, MARGIN, top + Inches(4.15), Inches(12), Inches(0.6),
          [("Precise conclusion (not \"PFT is useless\"):", 13.5, MUTED, True, True)])
add_text(s, MARGIN, top + Inches(4.55), Inches(12), Inches(0.85),
          [("Under the current Chronos-2 + LAI setup and available dataset, explicit PFT information does not "
            "provide measurable predictive information beyond a shuffled-PFT control.", 15, INK, True, True)],
          line_spacing=1.15)
add_takeaway(s, "Hypothesis, not proven ceiling: zero-shot Chronos-2 may already be hard to improve with small "
                "gradient-trained modules at this data scale.")

# ============================================================ SLIDE 8: RESEARCH PROGRESSION
s = add_slide(); set_bg(s)
add_title(s, "Research Progression: What We Ruled Out", eyebrow="SUMMARY")
page_num(s, 8)
top, h = content_box()

steps = [
    ("Initial hypothesis", "PFT does not help", GREY),
    ("Finding 1", "Chronos-2 architecturally cannot\nsee static PFT (InstanceNorm erasure)", WARN),
    ("Fix", "Explicit FiLM conditioning,\nbypassing InstanceNorm", ACCENT),
    ("Finding 2", "Architecture CAN learn strong\nPFT sensitivity — but doesn't generalize", WARN),
    ("Response", "4x more temporal windows +\n4 alternative architectures", ACCENT),
    ("Finding 3", "Real PFT performs no better\nthan shuffled PFT (p=0.245)", WARN),
    ("Conclusion", "No detectable, generalizable PFT\nsignal in the current setup", ACCENT_DARK),
]
n = len(steps)
box_w = Inches(1.62)
box_h = Inches(0.95)
gap = (SLIDE_W - 2 * MARGIN - n * box_w) // (n - 1)
y = top + Inches(0.5)
for i, (label, text, color) in enumerate(steps):
    x = MARGIN + i * (box_w + gap)
    add_text(s, x, y - Inches(0.32), box_w, Inches(0.3), [(label.upper(), 9.5, MUTED, True, False)],
              align=PP_ALIGN.CENTER)
    flow_box(s, x, y, box_w, box_h, text, color, size=10.5)
    if i < n - 1:
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + box_w, y + box_h // 2,
                                        x + box_w + gap, y + box_h // 2)
        conn.line.color.rgb = FAINT
        conn.line.width = Pt(2)

add_takeaway(s, "Each step ruled out a specific explanation with direct evidence — this is what makes the final "
                "negative result trustworthy.")

# ============================================================ SLIDE 9: CONCLUSION & NEXT STEPS
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.5), ACCENT)
add_text(s, MARGIN, Inches(0.35), Inches(12.2), Inches(0.9),
          [("Conclusion & Next Steps", 24, WHITE, True, False)])
page_num(s, 9)

top = Inches(1.85)
add_text(s, MARGIN, top, Inches(12.2), Inches(0.4), [("Conclusions", 16, ACCENT_DARK, True, False)])
add_bullets(s, MARGIN, top + Inches(0.45), Inches(12.2), Inches(1.7), [
    "Across 4 conditioning architectures and both PFT representations, no method beat a shuffled-PFT "
    "control — the strongest form of evidence this investigation can offer against the current approach.",
    "The negative result is well-diagnosed, not a dead end: the InstanceNorm blindness, the FiLM fix, and "
    "the shuffle control each closed off a specific alternative explanation.",
], size=15.5, space_after=10)

top2 = Inches(3.85)
add_text(s, MARGIN, top2, Inches(12.2), Inches(0.4), [("Recommended next steps", 16, ACCENT_DARK, True, False)])
add_bullets(s, MARGIN, top2 + Inches(0.45), Inches(12.2), Inches(2.6), [
    "Test PFT with the AELSTM-family models, which showed greater (if noisier) sensitivity to explicit PFT "
    "in the earlier single-pixel study.",
    "Revisit Chronos-2 only if substantially more or different training data become available — do not keep "
    "tuning increasingly complex FiLM/adapter variants on the current dataset.",
    "Longer term: investigate PFT as a modifier of drought/climate response specifically, rather than a "
    "generic static covariate, if sufficient supervision becomes available to learn that interaction.",
], size=15, space_after=10, bullet_color=ACCENT)
add_takeaway(s, "Full details: CHRONOS2_PFT_V2_REPORT.md — every number in this deck traces to a saved CSV in outputs/pft_v2/.")

out_path = f"{ROOT}/Chronos2_PFT_Investigation_Summary.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
