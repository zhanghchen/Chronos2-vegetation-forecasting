# -*- coding: utf-8 -*-
"""Comprehensive research summary deck, synthesizing all 6 research decks
produced 8/13-8/27: (1) Rolling-origin LOYO-CV + predictor sensitivity +
improved LoRA fine-tuning, (2) the 2012 leakage diagnostic, (3) PFT-aware
spatial transfer + cross-location drought training, (4) advanced PEFT
fine-tuning (6 methods vs. zero-shot), (5) the PFT architecture
investigation (InstanceNorm erasure -> FiLM fix -> multi-pixel), (6) the
final PFT-v2 shuffled-control result. All numbers are taken directly from
the text of the 6 source decks (themselves already verified against saved
CSVs in their own build scripts) - nothing here is re-derived or
invented."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image

ROOT = str(Path(__file__).resolve().parent.parent)
AELSTM_ROOT = "/home/deh25003/chronos-forecasting/AELSTM"
PFT_V2 = f"{ROOT}/outputs/pft_v2"
PFT_MULTI = f"{ROOT}/outputs/pft_multipixel"

INK = RGBColor(0x1C, 0x21, 0x19)
MUTED = RGBColor(0x5C, 0x63, 0x55)
FAINT = RGBColor(0x8B, 0x93, 0x82)
ACCENT = RGBColor(0x2F, 0x6F, 0x5E)
ACCENT_DARK = RGBColor(0x1E, 0x4A, 0x33)
ACCENT_TINT = RGBColor(0xE7, 0xF0, 0xE6)
WARN = RGBColor(0xB5, 0x65, 0x1D)
WARN_TINT = RGBColor(0xFB, 0xEE, 0xE0)
BLUE = RGBColor(0x3A, 0x6E, 0xA5)
BLUE_TINT = RGBColor(0xE7, 0xEE, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDE, 0xE2, 0xD6)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.25)
TAKEAWAY_H = Inches(0.78)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid(); bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or RULE; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, left, top, width, height, color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.12
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color; shp.line.width = Pt(1.0)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic; r.font.name = FONT
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=INK, space_after=9, bullet_color=ACCENT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.line_spacing = 1.06
        r0 = p.add_run(); r0.text = "▪  "; r0.font.size = Pt(size); r0.font.color.rgb = bullet_color; r0.font.name = FONT
        r1 = p.add_run(); r1.text = item; r1.font.size = Pt(size); r1.font.color.rgb = color; r1.font.name = FONT
    return box


def add_eyebrow(slide, text):
    add_text(slide, MARGIN, Inches(0.42), Inches(12), Inches(0.32), [(text, 13, ACCENT, True, False)])


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    add_text(slide, MARGIN, Inches(0.72), Inches(12.2), Inches(0.62), [(title, 23, INK, True, False)])
    add_rect(slide, MARGIN, Inches(1.38), Inches(1.0), Pt(3), ACCENT)


def add_takeaway(slide, text, tint=ACCENT_TINT, dark=ACCENT_DARK, bar=ACCENT):
    add_rect(slide, 0, TAKEAWAY_TOP, SLIDE_W, TAKEAWAY_H, tint)
    add_rect(slide, 0, TAKEAWAY_TOP, Inches(0.12), TAKEAWAY_H, bar)
    add_text(slide, Inches(0.45), TAKEAWAY_TOP, SLIDE_W - Inches(0.9), TAKEAWAY_H,
              [("TAKEAWAY   ", 12, dark, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.55), TAKEAWAY_TOP, SLIDE_W - Inches(2.0), TAKEAWAY_H,
              [(text, 14.5, dark, True, False)], anchor=MSO_ANCHOR.MIDDLE)


def page_num(slide, n):
    add_text(slide, SLIDE_W - Inches(0.7), Inches(0.42), Inches(0.4), Inches(0.3),
              [(str(n), 11, FAINT, False, False)], align=PP_ALIGN.RIGHT)


def content_box(has_takeaway=True):
    bottom = (TAKEAWAY_TOP - Inches(0.15)) if has_takeaway else (SLIDE_H - Inches(0.3))
    return CONTENT_TOP, bottom - CONTENT_TOP


def add_picture_fit(slide, path, box_left, box_top, box_w, box_h, align="center"):
    im = Image.open(path)
    ar = im.size[0] / im.size[1]
    box_ar = box_w / box_h
    if ar > box_ar:
        w = box_w; h = int(w / ar)
    else:
        h = box_h; w = int(h * ar)
    left = box_left + (box_w - w) // 2 if align == "center" else box_left
    top = box_top + (box_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    return left, top, w, h


def styled_table(slide, left, top, width, height, header, rows, col_weights, header_size=12, body_size=11.5, highlight_rows=()):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(header), left, top, width, height).table
    for c, w in enumerate(col_weights):
        tbl.columns[c].width = Emu(int(width * w))
    for ci, text in enumerate(header):
        cell = tbl.cell(0, ci); cell.text = text
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(header_size); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = FONT
        p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    for ri, row in enumerate(rows, start=1):
        hl = (ri - 1) in highlight_rows
        for ci, text in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = text
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_TINT if hl else (WHITE if ri % 2 else RGBColor(0xF5, 0xF6, 0xF1))
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(body_size); p.runs[0].font.bold = (ci == 0 or hl)
            p.runs[0].font.color.rgb = INK; p.runs[0].font.name = FONT
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    return tbl


def flow_box(slide, left, top, width, height, text, fill, text_color=WHITE, size=11):
    shp = add_rounded(slide, left, top, width, height, fill)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = text_color; r.font.bold = (i == 0); r.font.name = FONT
    return shp


def section_divider(title, subtitle, num, color=ACCENT):
    s = add_slide(); set_bg(s)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color)
    add_text(s, MARGIN, Inches(2.8), Inches(12.2), Inches(0.4), [(f"PART {num}", 15, WHITE, True, False)])
    add_text(s, MARGIN, Inches(3.2), Inches(12.2), Inches(1.1), [(title, 32, WHITE, True, False)])
    add_text(s, MARGIN, Inches(4.35), Inches(11.5), Inches(0.8), [(subtitle, 16, RGBColor(0xE0, 0xEC, 0xE6), False, True)], line_spacing=1.2)
    return s


# ============================================================ SLIDE 1: TITLE
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), ACCENT)
add_text(s, MARGIN, Inches(0.26), Inches(12.2), Inches(0.65), [("Vegetation LAI Forecasting: Comprehensive Research Summary", 24, WHITE, True, False)])
page_num(s, 1)
top, h = content_box()
add_text(s, MARGIN, top + Inches(0.1), Inches(11.8), Inches(0.4), [("SCOPE: AUGUST 13 – AUGUST 27, 2026", 13, ACCENT, True, False)])
add_bullets(s, MARGIN, top + Inches(0.65), Inches(11.8), Inches(3.6), [
    "Rolling-origin LOYO-CV, predictor sensitivity, and validation-based LoRA fine-tuning (8/13)",
    "2012 leakage diagnostic — does seeing the drought year in training explain the LOYO-CV failure? (8/14)",
    "PFT-aware spatial transfer and cross-location drought training for Georgia's 2012 drought (8/17)",
    "Advanced PEFT fine-tuning: 6 methods vs. zero-shot Chronos-2 (8/25)",
    "PFT architecture investigation: InstanceNorm erasure → FiLM conditioning → multi-pixel study → the decisive shuffled-PFT control (8/26–8/27)",
], size=17, space_after=13)
add_takeaway(s, "Same 3 core pixels, same model set, and raw-observation scoring underlie every study below.")

# ============================================================ SLIDE 2: SETUP RECAP
s = add_slide(); set_bg(s)
add_title(s, "Shared Setup Across All Studies", eyebrow="FOUNDATION")
page_num(s, 2)
top, h = content_box()
header = ["Pixel", "Location", "Land cover"]
rows = [
    ("low_amplitude", "37.53°N, 117.56°W", "Natural grassland"),
    ("high_amplitude_deciduous", "36.23°N, 84.47°W", "Broadleaf deciduous forest"),
    ("evergreen", "30.53°N, 82.43°W", "Needleleaf evergreen forest"),
]
styled_table(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(1.5), header, rows, [0.36, 0.32, 0.32], body_size=13)
add_bullets(s, MARGIN, top + Inches(1.85), SLIDE_W - 2 * MARGIN, Inches(2.6), [
    "8 AELSTM-family models: AELSTM, BiLSTM, LSTM, GRU, RNN, CNN, RF, SVM",
    "Chronos-2 (zero-shot and fine-tuned), using historical LAI + historical climate + known future climate",
    "7 gridMET predictors: tmmx, tmmn, pr, srad, vpd, sph, vs",
    "Baseline protocol: train 2000–2021, test 2022, scored against raw observed LAI",
], size=16.5, space_after=10)
add_takeaway(s, "All 6 studies below build on this identical foundation — differences are isolated to what's being tested.")

# ============================================================ PART 1: LOYO-CV
section_divider("Rolling-Origin LOYO-CV", "Is a single 2022 test year enough to trust a model comparison?", 1)

s = add_slide(); set_bg(s)
add_title(s, "Design: 11 Folds, Fixed 12-Year Window", eyebrow="PART 1 · LOYO-CV")
page_num(s, 4)
top, h = content_box()
add_bullets(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(1.7), [
    "2022 turned out not to be an anomalous year at any of the 3 pixels — a single fixed test year gives "
    "no evidence about how models behave when a real anomaly occurs.",
    "Fixed 12-year training window immediately before each held-out year (not classic all-other-years LOYO, "
    "not an expanding window) — isolates genuine year-to-year difficulty from training-set-size effects.",
], size=16, space_after=10)
windows = [("2000–2011", "2012"), ("2001–2012", "2013"), ("...", "..."), ("2010–2021", "2022")]
x = MARGIN
box_w = Inches(2.75)
for i, (train, test) in enumerate(windows):
    flow_box(s, x, top + Inches(2.0), box_w, Inches(0.9), f"train {train}\ntest {test}\nfold {i+1 if i<3 else 11}", ACCENT if i < 3 else ACCENT_DARK, size=12)
    x += box_w + Inches(0.2)
add_takeaway(s, "11 folds (2012–2022) × same 10 methods × same 3 pixels, scored against raw observed LAI.")

s = add_slide(); set_bg(s)
add_title(s, "Findings: Stability, Reliable Models, and Two Outlier Folds", eyebrow="PART 1 · LOYO-CV")
page_num(s, 5)
top, h = content_box()
add_bullets(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.2), [
    "high_amplitude_deciduous: all 10 methods score R² 0.89–0.99 across 11 years — stable regardless of "
    "method. low_amplitude and evergreen show a much wider spread, each with a severe outlier fold.",
    "RF (mean rank 3.39) and Chronos-2 zero-shot (3.70) are top-3 on ~6 of every 10 folds and rarely "
    "worst — the most dependable models across all 33 folds. LoRA fine-tuning is the single most volatile "
    "method; AELSTM and CNN are the least reliable overall.",
    "Two outlier folds, two different causes: evergreen/2012 (z=−3.0) is a real, sustained drought; "
    "low_amplitude/2018 (z=−2.9) is not a climate anomaly at all, just an erratic year at an already "
    "low-signal pixel. No year was difficult at all 3 pixels simultaneously.",
], size=16.5, space_after=13)
add_takeaway(s, "A fixed 2022-only test could never have revealed either failure mode — report medians alongside means.")

# ============================================================ PART 1b: LEAKAGE DIAGNOSTIC
s = add_slide(); set_bg(s)
add_title(s, "2012 Leakage Diagnostic — Distribution Shift or Model Limitation?", eyebrow="FOLLOW-UP")
page_num(s, 6)
top, h = content_box()
add_rect(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.75), BLUE_TINT)
add_text(s, MARGIN + Inches(0.25), top + Inches(0.1), SLIDE_W - 3 * MARGIN, Inches(0.55),
          [("Question: \"What happens if you predict 2012 LAI using a model that has already seen 2012 during training?\"", 14.5, BLUE, True, True)],
          anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, MARGIN, top + Inches(1.0), SLIDE_W - 2 * MARGIN, Inches(3.2), [
    "AELSTM-family (8 models, full refit): mean RMSE ↓ ~46%, mean R² ↑ +5.17 once 2012 is included in "
    "fitting — consistent across very different architectures. RF recovers the most (R²=+0.59, the only "
    "model to cross into positive territory); the other 7 improve substantially but stay negative.",
    "Chronos-2 (LoRA, the only way to let it \"see\" 2012): RMSE ↓ 59.9%, R² −8.29 → −0.50, Pearson r "
    "0.13 → 0.65, ACC −0.63 → 0.88 — the same large-recovery-but-still-short-of-good-fit pattern.",
    "Answer: the dominant cause is distribution shift (2012 was a real drought unlike anything in the "
    "training history) — but model-specific limitations still remain, since most models stay negative "
    "even after seeing 2012.",
], size=15.5, space_after=10)
add_takeaway(s, "Distribution shift explains most, but not all, of the original LOYO-CV failure at evergreen/2012.")

# ============================================================ PART 2: PREDICTOR SENSITIVITY
s = add_slide(); set_bg(s)
add_title(s, "Predictor Sensitivity: Design and Findings", eyebrow="PART 2")
page_num(s, 7)
top, h = content_box()
add_bullets(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.4), [
    "Design: leave-one-predictor-out (7 configs) + 4 grouped ablations (selection rule fixed before Phase "
    "1 ran) — a small, pre-registered design instead of an intractable 2⁷=128-subset search.",
    "Predictor importance is model- and pixel-specific: the same removal helps one model and hurts another "
    "at the same pixel; high_amplitude_deciduous is largely insensitive to any single predictor.",
    "srad (solar radiation) is the only predictor that helps at all 3 pixels — the single most consistently "
    "important variable. Temperature is critical at low_amplitude but unimportant elsewhere.",
    "Dropping the 2 least-important predictors (vs, pr) improves mean R² at all 3 pixels simultaneously — "
    "a leaner 5-predictor set (tmmx, tmmn, srad, vpd, sph) is a validated, better default.",
], size=16, space_after=11)
add_takeaway(s, "No single predictor-importance ranking applies universally — but srad is consistently valuable, vs/pr are redundant.")

# ============================================================ PART 3: IMPROVED LoRA
s = add_slide(); set_bg(s)
add_title(s, "Improved LoRA Fine-Tuning: Validation Recovers Most of the Loss", eyebrow="PART 3")
page_num(s, 8)
top, h = content_box()
add_bullets(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(1.3), [
    "The original fine-tune had no safety net: fit() was called without validation_inputs, trained a fixed "
    "1000 steps, and kept whatever weights existed at that instant — no way to detect overfitting.",
], size=15.5, space_after=8)
header2 = ["Pixel", "Zero-shot", "Original LoRA", "Improved LoRA"]
rows2 = [
    ("low_amplitude", "0.542", "0.379", "0.500"),
    ("high_amplitude_deciduous", "0.965", "0.952", "0.942"),
    ("evergreen", "0.832", "0.712", "0.829"),
]
styled_table(s, MARGIN, top + Inches(1.5), SLIDE_W - 2 * MARGIN, Inches(1.5), header2, rows2, [0.34, 0.22, 0.22, 0.22], body_size=13)
add_bullets(s, MARGIN, top + Inches(3.15), SLIDE_W - 2 * MARGIN, Inches(1.4), [
    "Chronological validation (2020 & 2021 folds, 2022 untouched) + a small search (lr × rank) with early "
    "stopping: gap to zero-shot shrinks from −0.163 to −0.042 at low_amplitude, and −0.120 to −0.003 at "
    "evergreen — but improved LoRA still does not outright beat zero-shot on any of the 3 pixels.",
], size=15.5, space_after=8)
add_takeaway(s, "Most of the original degradation was a training-setup artifact, not fundamental to LoRA itself.")

# ============================================================ PART 4: PFT-AWARE SPATIAL TRANSFER
section_divider("PFT-Aware Spatial Transfer", "Can other drought-affected locations help predict Georgia's 2012 drought?", 4)

s = add_slide(); set_bg(s)
add_title(s, "First Attempt: PFT-Match Alone Wasn't Enough", eyebrow="PART 4 · SPATIAL TRANSFER")
page_num(s, 10)
top, h = content_box()
add_bullets(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.4), [
    "10 needleleaf-evergreen pixels (same PFT as Georgia, fraction ≥70%) selected across 9 US regions — "
    "but checking each candidate's 2000–2011 climatology against actual 2012 LAI showed only 2 of 9 "
    "carried a real drought anomaly in the same direction as Georgia.",
    "Refined design: A. Local (Georgia 2000–2011 only), B. Single-source (best matched pixel), "
    "C. Multi-source (5 pixels: 2 informative + 3 diverse), D. Diagnostic reference (Georgia's own "
    "2012, leaked).",
    "Mean R²: A = −6.56, B = −2.98, C = −1.18, D (leaked reference) = −1.39 — C (no Georgia 2012 data at "
    "all) performs essentially as well as D (Georgia's own 2012 leaked into training).",
], size=16, space_after=12)
add_takeaway(s, "Select source pixels by evidence of the actual anomaly, not by vegetation-type label alone.")

s = add_slide(); set_bg(s)
add_title(s, "Refined: Cross-Location Drought Training", eyebrow="PART 4 · REFINED DESIGN")
page_num(s, 11)
top, h = content_box()
header3 = ["Condition", "Mean R²", "Mean RMSE"]
rows3 = [("A. Local baseline", "-6.56", "1.58"), ("B. Single source", "-2.31", "0.95"), ("C. Multi-source (5)", "-2.32", "1.05")]
styled_table(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(1.4), header3, rows3, [0.5, 0.25, 0.25], body_size=14)
add_bullets(s, MARGIN, top + Inches(1.7), SLIDE_W - 2 * MARGIN, Inches(2.8), [
    "Selection changed from PFT-label-first to drought-evidence-first: rank candidates by their own "
    "2012 anomaly strength (mean z-score, peak-LAI reduction), with PFT kept only as a constraint — "
    "found source pixels with a 5–10x stronger 2012 anomaly than the earlier PFT-only screen's best.",
    "Both B and C cut the local baseline's damage dramatically, without ever using Georgia's own 2012 data.",
    "Secondary analysis — number of source pixels (N=1,2,3,5,10): R² = −2.31, −4.00, −2.78, −2.32, −3.27 "
    "— not monotonic. Pooling all 10 candidates performs worse than a well-chosen 5.",
], size=15.5, space_after=9)
add_takeaway(s, "Selecting genuinely informative source pixels matters more than simply increasing pixel count.")

# ============================================================ PART 5: ADVANCED FINE-TUNING
section_divider("Advanced PEFT Fine-Tuning", "Can a more advanced fine-tuning strategy beat zero-shot Chronos-2?", 5)

s = add_slide(); set_bg(s)
add_title(s, "Six PEFT Methods vs. Zero-Shot", eyebrow="PART 5 · ADVANCED FINE-TUNING")
page_num(s, 13)
top, h = content_box()
if Path(f"{ROOT}/outputs/advanced_finetuning/r2_by_pixel_all_methods.png").exists():
    add_picture_fit(s, f"{ROOT}/outputs/advanced_finetuning/r2_by_pixel_all_methods.png", MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.1))
add_takeaway(s, "DoRA, VeRA, IA3, LN-Tuning, BitFit, partial-last-block — literature-grounded, mechanistically diverse methods.")

s = add_slide(); set_bg(s)
add_title(s, "No Method Beats Zero-Shot; Capacity Doesn't Predict Performance", eyebrow="PART 5 · ADVANCED FINE-TUNING")
page_num(s, 14)
top, h = content_box()
header4 = ["Method", "Params", "Mean R²", "vs. Zero-shot"]
rows4 = [
    ("Zero-shot", "—", "0.780", "—"),
    ("IA3", "74K", "0.769", "-0.011"),
    ("Partial (last block)", "13.1M", "0.766", "-0.014"),
    ("DoRA", "1.3M", "0.761", "-0.019"),
    ("BitFit", "8.4K", "0.683", "-0.097"),
    ("Original LoRA", "—", "0.681", "-0.099"),
]
styled_table(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(2.6), header4, rows4, [0.32, 0.18, 0.22, 0.28], body_size=13, highlight_rows=(0,))
add_bullets(s, MARGIN, top + Inches(2.9), SLIDE_W - 2 * MARGIN, Inches(1.6), [
    "Trainable parameter count does not predict performance: IA3 (74K params) beats partial fine-tuning "
    "(13.1M params); BitFit (8.4K, smallest) is the worst performer overall.",
    "BitFit's winning config had the BEST validation loss of all 24 evergreen configs, yet the WORST "
    "test R² — a genuine failure mode, not classic overfitting (train/val curves looked normal).",
], size=14.5, space_after=7)
add_takeaway(s, "Zero-shot Chronos-2 is already a strong baseline — no PEFT method beat it on this dataset.")

# ============================================================ PART 6: PFT INVESTIGATION
section_divider("The PFT Investigation", "Does vegetation-composition information improve Chronos-2's LAI forecasts?", 6)

s = add_slide(); set_bg(s)
add_title(s, "Diagnosis: Chronos-2 Was Architecturally Blind to Static PFT", eyebrow="PART 6 · PFT")
page_num(s, 16)
top, h = content_box()
diagram_top = top + Inches(0.3)
box_w, box_h = Inches(2.6), Inches(0.75)
gap = Inches(0.55)
x0 = MARGIN + Inches(0.3)
flow_box(s, x0, diagram_top, box_w, box_h, "PFT = 0.65\n(constant every step)", MUTED, size=12.5)
flow_box(s, x0 + box_w + gap, diagram_top, box_w, box_h, "InstanceNorm\n(per-series, std=0)", WARN, size=12.5)
flow_box(s, x0 + 2 * (box_w + gap), diagram_top, box_w, box_h, "→ [0, 0, 0, ...]\nsignal erased", ACCENT_DARK, size=12.5)
for i in range(2):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x0 + (i + 1) * box_w + i * gap, diagram_top + box_h // 2,
                                    x0 + (i + 1) * box_w + (i + 1) * gap, diagram_top + box_h // 2)
    conn.line.color.rgb = FAINT; conn.line.width = Pt(2.5)
add_bullets(s, MARGIN, diagram_top + Inches(1.15), SLIDE_W - 2 * MARGIN, Inches(2.5), [
    "Verified empirically: perturbing PFT across 5 synthetic compositions changed the prediction by "
    "exactly 0.00000000 — not a small effect, zero.",
    "Fix: a small trainable FiLM-conditioning module (~12.6K–100.5K params, 0.01–0.08% of the model) "
    "injects PFT directly into the model's internal representation, bypassing InstanceNorm entirely — "
    "the 119.5M-parameter base model stays completely frozen.",
    "Confirmed working: a smoke test showed a large, structured response to synthetic PFT perturbation.",
], size=15.5, space_after=9)
add_takeaway(s, "This was a representation problem, not evidence that PFT is biologically irrelevant.")

s = add_slide(); set_bg(s)
add_title(s, "Multi-Pixel Study: Overfitting, Even With 4x More Data", eyebrow="PART 6 · PFT")
page_num(s, 17)
top, h = content_box()
if Path(f"{PFT_V2}/deck_figures/overfitting_comparison.png").exists():
    add_picture_fit(s, f"{PFT_V2}/deck_figures/overfitting_comparison.png", MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(4.0))
add_takeaway(s, "70 PFT-diverse pixels pooled; 3 of 4 conditioning architectures still overfit within one training step.")

# ============================================================ SLIDE: KEY FINDING
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), WARN)
add_text(s, MARGIN, Inches(0.28), Inches(12.2), Inches(0.6), [("Key Finding: Real PFT ≈ Shuffled PFT", 25, WHITE, True, False)])
page_num(s, 18)
top, h = content_box()
if Path(f"{PFT_V2}/deck_figures/real_vs_shuffled_bar.png").exists():
    add_picture_fit(s, f"{PFT_V2}/deck_figures/real_vs_shuffled_bar.png", MARGIN, top, SLIDE_W - 2 * MARGIN, h - Inches(0.05))
add_takeaway(s, "A model trained on randomly shuffled PFT performs just as well as one trained on real PFT — p=0.245, not significant.",
             tint=WARN_TINT, dark=WARN, bar=WARN)

s = add_slide(); set_bg(s)
add_title(s, "What This Means", eyebrow="PART 6 · INTERPRETATION")
page_num(s, 19)
top, h = content_box()
add_rect(s, MARGIN, top + Inches(0.05), Inches(11.6), Inches(0.9), ACCENT_TINT)
add_text(s, MARGIN + Inches(0.3), top + Inches(0.17), Inches(11), Inches(0.7),
          [("No PFT  ≈  Real PFT  ≈  Shuffled PFT", 19, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, MARGIN, top + Inches(1.15), SLIDE_W - 2 * MARGIN, Inches(2.4), [
    "The small ΔR² most likely comes from adding trainable conditioning capacity itself, not from the "
    "biological information the PFT vector carries.",
    "Mixed-vs-pure entropy, seasonal-phase RMSE, and per-dominant-PFT-class breakdowns all reach the "
    "same conclusion once compared against the shuffled control.",
    "Precise conclusion (not \"PFT is useless\"): under the current Chronos-2 + LAI setup and available "
    "dataset, explicit PFT information does not provide a measurable, generalizable predictive benefit "
    "beyond a shuffled-PFT control.",
], size=15.5, space_after=10)
add_takeaway(s, "This is a hypothesis about the current setup and data scale, not a claim of a theoretical ceiling.")

# ============================================================ SYNTHESIS
s = add_slide(); set_bg(s)
add_title(s, "Cross-Cutting Synthesis: Two Threads, One Observation", eyebrow="SYNTHESIS")
page_num(s, 20)
top, h = content_box()
add_bullets(s, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(3.4), [
    "Fine-tuning thread: LOYO-CV showed Chronos-2 zero-shot is among the most consistently strong models "
    "across 33 folds; the improved-LoRA study and the 6-method advanced PEFT comparison both found no "
    "fine-tuning strategy consistently beats zero-shot, regardless of trainable parameter count.",
    "PFT thread: the architecture was fixed so it CAN receive PFT information, and can express a strong "
    "response to it — but a properly validated model, and the shuffled-PFT control specifically, found no "
    "evidence it learns anything real from it under the current data scale.",
    "Both threads converge on the same empirical pattern: zero-shot Chronos-2 is difficult to improve on "
    "this dataset with relatively small amounts of gradient-trained adaptation — whether the adaptation "
    "targets general fine-tuning or PFT-specific conditioning.",
], size=16, space_after=12)
add_takeaway(s, "An empirical observation about this dataset/setup, not evidence of a theoretical performance ceiling.")

# ============================================================ CONCLUSIONS
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.5), ACCENT)
add_text(s, MARGIN, Inches(0.35), Inches(12.2), Inches(0.9), [("Overall Conclusions & Next Steps", 24, WHITE, True, False)])
page_num(s, 21)
top = Inches(1.85)
add_text(s, MARGIN, top, Inches(12.2), Inches(0.4), [("What We Learned", 16, ACCENT_DARK, True, False)])
add_bullets(s, MARGIN, top + Inches(0.45), Inches(12.2), Inches(2.3), [
    "Rolling-origin LOYO-CV is essential — single-year evaluation misses failure modes entirely (RF and "
    "Chronos-2 zero-shot are the most consistently strong models across 33 folds).",
    "Source-pixel quality (evidence of the actual phenomenon) matters more than PFT-label match or pixel "
    "count for spatial transfer.",
    "Across every fine-tuning method tested and the PFT-conditioning architecture, zero-shot Chronos-2 has "
    "been difficult to consistently improve on with the data and methods tried so far.",
], size=15.5, space_after=9)

top2 = Inches(4.55)
add_text(s, MARGIN, top2, Inches(12.2), Inches(0.4), [("Next Steps", 16, ACCENT_DARK, True, False)])
add_bullets(s, MARGIN, top2 + Inches(0.45), Inches(12.2), Inches(2.0), [
    "Extend LOYO-CV / predictor ablation to the remaining selected pixels; adopt the leaner 5-predictor set.",
    "Test PFT with AELSTM-family models, which showed greater (if noisier) sensitivity to explicit PFT.",
    "Revisit Chronos-2 adaptation (fine-tuning and PFT alike) if substantially more training data or "
    "stronger supervision becomes available.",
], size=15.5, space_after=9, bullet_color=ACCENT)
add_takeaway(s, "Full detail for every study lives in its own report/deck — this deck is the cross-study summary.")

out_path = f"{ROOT}/Vegetation_Forecasting_Comprehensive_Summary.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
