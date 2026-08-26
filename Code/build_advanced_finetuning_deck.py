# -*- coding: utf-8 -*-
"""Builds the advanced Chronos-2 fine-tuning summary deck (8 slides).
Reuses the exact visual identity from the project's other decks
(build_cross_location_deck.py / build_pft_spatial_deck.py). All numbers/
figures pulled directly from CHRONOS2_ADVANCED_FINETUNING_REPORT.md and
outputs/advanced_finetuning/ (commit 39d715d) - nothing here is invented."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = str(Path(__file__).resolve().parent.parent)
OUT_DATA = f"{ROOT}/outputs/advanced_finetuning"

INK = RGBColor(0x1C, 0x21, 0x19)
MUTED = RGBColor(0x5C, 0x63, 0x55)
FAINT = RGBColor(0x8B, 0x93, 0x82)
ACCENT = RGBColor(0x2F, 0x6F, 0x4E)
ACCENT_DARK = RGBColor(0x1E, 0x4A, 0x33)
ACCENT_TINT = RGBColor(0xE7, 0xF0, 0xE6)
WARN = RGBColor(0xB5, 0x65, 0x1D)
WARN_TINT = RGBColor(0xFB, 0xEE, 0xE0)
BLUE = RGBColor(0x3A, 0x6E, 0xA5)
BLUE_TINT = RGBColor(0xE7, 0xEE, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDE, 0xE2, 0xD6)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.25)
TAKEAWAY_H = Inches(0.78)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def emu(v):
    return Emu(int(v))


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = RULE
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
    return box


def add_bullets(slide, left, top, width, height, items, size=17, color=INK, space_after=10,
                 bullet_color=ACCENT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        r0 = p.add_run()
        r0.text = "▪  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        r0.font.name = FONT
        r1 = p.add_run()
        r1.text = item
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.name = FONT
    return box


def add_eyebrow(slide, text):
    add_text(slide, MARGIN, Inches(0.42), Inches(12), Inches(0.32), [(text, 13, ACCENT, True, False)])


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    add_text(slide, MARGIN, Inches(0.72), Inches(12.2), Inches(0.62), [(title, 25, INK, True, False)])
    add_rect(slide, MARGIN, Inches(1.38), Inches(1.0), Pt(3), ACCENT)


def add_takeaway(slide, text):
    add_rect(slide, 0, TAKEAWAY_TOP, SLIDE_W, TAKEAWAY_H, ACCENT_TINT)
    add_rect(slide, 0, TAKEAWAY_TOP, Inches(0.12), TAKEAWAY_H, ACCENT)
    add_text(slide, Inches(0.45), TAKEAWAY_TOP, SLIDE_W - Inches(0.9), TAKEAWAY_H,
              [("TAKEAWAY   ", 12, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.55), TAKEAWAY_TOP, SLIDE_W - Inches(2.0), TAKEAWAY_H,
              [(text, 15.5, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)


def page_num(slide, n):
    add_text(slide, SLIDE_W - Inches(0.7), Inches(0.42), Inches(0.4), Inches(0.3),
              [(str(n), 11, FAINT, False, False)], align=PP_ALIGN.RIGHT)


def add_picture_fit(slide, path, box_left, box_top, box_w, box_h, align="center"):
    im = Image.open(path)
    ar = im.size[0] / im.size[1]
    box_ar = box_w / box_h
    if ar > box_ar:
        w = box_w
        h = int(w / ar)
    else:
        h = box_h
        w = int(h * ar)
    left = box_left + (box_w - w) // 2 if align == "center" else box_left
    top = box_top + (box_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    return left, top, w, h


def content_box(has_takeaway=True):
    bottom = (TAKEAWAY_TOP - Inches(0.15)) if has_takeaway else (SLIDE_H - Inches(0.3))
    return CONTENT_TOP, bottom - CONTENT_TOP


def styled_table(slide, left, top, width, height, header, rows, col_weights,
                  header_size=12, body_size=11.5, highlight_rows=()):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(header), left, top, width, height).table
    for c, w in enumerate(col_weights):
        tbl.columns[c].width = Emu(int(width * w))
    for ci, text in enumerate(header):
        cell = tbl.cell(0, ci)
        cell.text = text
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(header_size); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = FONT
        p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    for ri, row in enumerate(rows, start=1):
        hl = (ri - 1) in highlight_rows
        for ci, text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_TINT if hl else (WHITE if ri % 2 else RGBColor(0xF5, 0xF6, 0xF1))
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(body_size); p.runs[0].font.bold = (ci == 0 or hl)
            p.runs[0].font.color.rgb = INK; p.runs[0].font.name = FONT
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    return tbl


# ============================================================ SLIDE 1: TITLE / RESEARCH QUESTION
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), ACCENT)
add_text(s, MARGIN, Inches(0.3), Inches(12), Inches(0.6),
          [("Advanced Chronos-2 Fine-Tuning: Can Better PEFT Beat Zero-Shot?", 24, WHITE, True, False)])
page_num(s, 1)

top, h = content_box()
add_text(s, MARGIN, top + Inches(0.1), Inches(12.2), Inches(0.4), [("RESEARCH QUESTION", 13, ACCENT, True, False)])
add_rect(s, MARGIN, top + Inches(0.55), Inches(0.1), Inches(1.1), ACCENT)
add_text(s, MARGIN + Inches(0.35), top + Inches(0.5), Inches(11.8), Inches(1.2),
          [("Can a more advanced fine-tuning strategy consistently outperform Chronos-2 zero-shot "
            "for vegetation LAI forecasting under limited-data conditions?", 20, INK, True, True)],
          line_spacing=1.15)
add_bullets(s, MARGIN, top + Inches(2.05), Inches(12.2), Inches(2.2), [
    "Standard LoRA fine-tuning (original and validation-selected) did not consistently beat "
    "zero-shot in our earlier experiments.",
    "This study tests 6 additional, more recent parameter-efficient fine-tuning (PEFT) methods, "
    "selected from a literature review of ICLR/ICML/NeurIPS 2022-2025.",
    "Same 3 pixels, same fair validation-only protocol, same pretrained starting weights for "
    "every method — an equal-sized hyperparameter budget throughout.",
], size=16.5, space_after=12)
add_takeaway(s, "6 mechanistically distinct methods, 90 fine-tuning runs, one fair protocol — does anything beat zero-shot?")

# ============================================================ SLIDE 2: LITERATURE REVIEW
s = add_slide(); set_bg(s)
add_title(s, "Literature Review & Method Selection")
page_num(s, 2)
top, h = content_box()
header = ["Method", "Venue", "Core idea", "Trained params"]
rows = [
    ("DoRA", "ICML'24 Oral", "Decompose weight into magnitude + direction; LoRA-update direction", "~1.3M"),
    ("VeRA", "ICLR'24", "One frozen shared random low-rank pair; learn only small scale vectors", "~99K"),
    ("IA3", "NeurIPS'22", "Learned multiplicative rescaling of K/V/FFN activations", "~74K"),
    ("LN-Tuning", "Chronos precedent (2409.11302)", "Unfreeze only normalization affine parameters", "28K"),
    ("BitFit", "ACL'22", "Unfreeze only bias terms", "8.4K"),
    ("Partial (last block)", "Classical", "Full-weight unfreeze of only the last transformer block", "~13.1M"),
]
styled_table(s, MARGIN, top + Inches(0.05), SLIDE_W - 2 * MARGIN, Inches(2.9), header, rows,
             [0.16, 0.19, 0.5, 0.15], body_size=12)
add_bullets(s, MARGIN, top + Inches(3.15), SLIDE_W - 2 * MARGIN, Inches(1.6), [
    "Chosen to span mechanistically distinct update types (not 5 near-identical LoRA variants) and a "
    "capacity spectrum from 8.4K to 13.1M trainable parameters.",
    "Excluded after investigation: AdaLoRA (needs more gradient signal than our tiny datasets give), "
    "FourierFT (overlaps with VeRA), prompt/prefix tuning (real compatibility risk — Chronos-2 isn't a "
    "token-sequence LM), PETSA/TimesFM in-context tuning (break comparability with our fixed protocol).",
], size=14.5, space_after=8)
add_takeaway(s, "Selection grounded in real compatibility checks against Chronos-2's actual source code, not just fame.")

# ============================================================ SLIDE 3: EXPERIMENTAL DESIGN
s = add_slide(); set_bg(s)
add_title(s, "Experimental Design")
page_num(s, 3)
top, h = content_box()
add_bullets(s, MARGIN, top + Inches(0.05), Inches(12.2), Inches(2.0), [
    "Every method starts from the identical pretrained Chronos-2 weights — verified by reading the "
    "actual fit() source; Chronos2Trainer has zero LoRA-specific logic, so any peft.PeftConfig can be "
    "substituted via a faithful reimplementation of fit()'s internals.",
    "Same validated protocol as the improved-LoRA study: chronological validation folds (2020, 2021), "
    "early stopping, final refit on 2000-2021, evaluated once on real 2022 — test data never touches "
    "hyperparameter selection.",
    "Equal-sized 4-config search budget per method per pixel — no method received more tuning effort "
    "than another.",
], size=17, space_after=14)

grid_top = top + Inches(2.6)
add_rect(s, MARGIN, grid_top, SLIDE_W - 2 * MARGIN, Inches(1.5), ACCENT_TINT)
add_text(s, MARGIN + Inches(0.3), grid_top + Inches(0.15), SLIDE_W - 2 * MARGIN - Inches(0.6), Inches(1.2),
          [("DoRA: lr∈{1e-5,1e-4} × rank∈{8,16}     VeRA: lr∈{1e-4,1e-3} × rank∈{256,1024}", 14, ACCENT_DARK, True, False),
           ("IA3 / LN-Tuning / BitFit: lr∈{1e-4,1e-3,1e-2,1e-1}     Partial: lr∈{1e-6,1e-5,1e-4,1e-3}", 14, ACCENT_DARK, True, False),
           ("→ 72 search runs + 18 final refits = 90 fine-tuning runs, ~85 minutes on one A100", 14, ACCENT_DARK, True, True)],
          line_spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)
add_takeaway(s, "Baselines (zero-shot, original LoRA, improved LoRA) reused directly — never rerun.")

# ============================================================ SLIDE 4: RESULTS BY PIXEL
s = add_slide(); set_bg(s)
add_title(s, "Results: R² by Pixel, All 9 Methods")
page_num(s, 4)
top, h = content_box()
add_picture_fit(s, f"{OUT_DATA}/r2_by_pixel_all_methods.png", MARGIN, top, SLIDE_W - 2 * MARGIN, h - Inches(0.1))
add_takeaway(s, "Zero-shot (dashed line) sits at or near the top of every panel — no method consistently clears it.")

# ============================================================ SLIDE 5: FINAL RANKING TABLE
s = add_slide(); set_bg(s)
add_title(s, "Final Ranking Table")
page_num(s, 5)
top, h = content_box()
header2 = ["Method", "Params", "Mean R²", "vs. Zero-shot", "Pixels beating ZS"]
rows2 = [
    ("Zero-shot", "—", "0.780", "—", "—"),
    ("IA3", "74K", "0.769", "-0.011", "0/3"),
    ("Partial (last block)", "13.1M", "0.766", "-0.014", "0/3"),
    ("DoRA", "1.3M", "0.761", "-0.019", "1/3"),
    ("LN-Tuning", "28K", "0.760", "-0.020", "0/3"),
    ("Improved LoRA", "—", "0.757", "-0.023", "0/3"),
    ("VeRA", "99K", "0.755", "-0.025", "0/3"),
    ("BitFit", "8.4K", "0.683", "-0.097", "0/3"),
    ("Original LoRA", "—", "0.681", "-0.099", "0/3"),
]
styled_table(s, MARGIN, top + Inches(0.1), SLIDE_W - 2 * MARGIN, Inches(4.0), header2, rows2,
             [0.28, 0.14, 0.16, 0.2, 0.22], header_size=13, body_size=13, highlight_rows=(0,))
add_takeaway(s, "Only one method/pixel combination ever beats zero-shot — DoRA on evergreen, by +0.0026.")

# ============================================================ SLIDE 6: CAPACITY VS PERFORMANCE
s = add_slide(); set_bg(s)
add_title(s, "Does More Capacity Help?")
page_num(s, 6)
top, h = content_box()
add_picture_fit(s, f"{OUT_DATA}/params_vs_r2.png", MARGIN, top, Inches(7.6), h - Inches(0.1), align="left")
add_bullets(s, MARGIN + Inches(7.9), top + Inches(0.3), Inches(4.5), h, [
    "Overall correlation between log(trainable params) and R² is weak: only +0.10 pooled across "
    "all pixels.",
    "Direction flips by pixel: +0.59 (evergreen), +0.45 (low_amplitude), but −0.24 "
    "(high_amplitude_deciduous) — more capacity mildly hurts where zero-shot is already near-ceiling.",
    "Which pixel matters far more than how much capacity a method has.",
], size=15, space_after=12)
add_takeaway(s, "Parameter count is a weak, pixel-dependent predictor — not a reliable guide to which method to pick.")

# ============================================================ SLIDE 7: FAILURE MODE
s = add_slide(); set_bg(s)
add_title(s, "Failure Mode: When Validation Loss Lies")
page_num(s, 7)
top, h = content_box()
add_picture_fit(s, f"{OUT_DATA}/validation_curves_bitfit_vs_dora_evergreen.png", MARGIN, top,
                  SLIDE_W - 2 * MARGIN, Inches(4.0))
add_bullets(s, MARGIN, top + Inches(4.15), SLIDE_W - 2 * MARGIN, Inches(1.0), [
    "BitFit's selected config had the BEST validation loss of all 24 evergreen configs tested (0.562, "
    "beating even DoRA's 0.615) — yet the WORST test R² (0.587 vs. DoRA's 0.835).",
    "Not classic overfitting (train/val curves look normal) — BitFit's only lever (a patch-embedding "
    "bias shift) helped the 2020/2021 folds but pointed the wrong way for 2022.",
], size=14.5, space_after=6)
add_takeaway(s, "Chronos-2's attention/FFN layers have NO bias terms — BitFit can only touch 0.007% of the model.")

# ============================================================ SLIDE 8: CONCLUSION
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.5), ACCENT)
add_text(s, MARGIN, Inches(0.28), Inches(12.2), Inches(0.9),
          [("No — zero-shot Chronos-2 remains the best method overall for this task.", 22, WHITE, True, False)],
          line_spacing=1.1)
page_num(s, 8)

top = Inches(1.85)
add_bullets(s, MARGIN, top, Inches(12.2), Inches(2.2), [
    "Across 6 mechanistically distinct methods (additive low-rank, shared-random-projection, "
    "multiplicative rescaling, bias-only, normalization-only, partial full-weight) and 8K-13M "
    "trainable parameters, not one consistently beat zero-shot across the 3 pixels.",
    "This corroborates independent literature evidence that Chronos-2 fine-tuning specifically "
    "struggles on small datasets, and extends our project's LoRA-specific finding: the result is not "
    "about LoRA's particular mechanism — it looks like a property of fine-tuning Chronos-2 at all on "
    "data this small.",
], size=16.5, space_after=12)

top2 = Inches(4.5)
add_text(s, MARGIN, top2, Inches(12.2), Inches(0.4), [("Practical recommendation", 16, ACCENT_DARK, True, False)])
add_bullets(s, MARGIN, top2 + Inches(0.45), Inches(12.2), Inches(1.8), [
    "Keep zero-shot Chronos-2 as the default for this task.",
    "If fine-tuning is needed for a structurally different domain, IA3 or last-block-only partial "
    "fine-tuning are the safest choices tested — they lose the least relative to zero-shot with no "
    "catastrophic single-pixel failure.",
], size=15.5, space_after=10, bullet_color=ACCENT)
add_takeaway(s, "Method diversity ruled out LoRA's mechanism as the cause — the ceiling appears to be the data, not the method.")

out_path = f"{ROOT}/Chronos2_Advanced_Finetuning_Summary.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
