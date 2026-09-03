# -*- coding: utf-8 -*-
"""Weekly progress-update deck: ERA5 meteorological data preparation for
the global Chronos-2 LAI forecasting experiment. Companion to the weekly
group-update email. Every fact here is taken directly from
experiments/global_era5_chronos/ (era5_source.py's module docstring,
select_global_pixels.py's output CSV, run_era5_fetch_supervised.sh and
its log, and common_pipeline.py) - nothing is invented or re-derived."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = str(Path(__file__).resolve().parent.parent)

INK = RGBColor(0x1C, 0x21, 0x19)
MUTED = RGBColor(0x5C, 0x63, 0x55)
FAINT = RGBColor(0x8B, 0x93, 0x82)
ACCENT = RGBColor(0x2F, 0x6F, 0x5E)
ACCENT_DARK = RGBColor(0x1E, 0x4A, 0x33)
ACCENT_TINT = RGBColor(0xE7, 0xF0, 0xE6)
WARN = RGBColor(0xB5, 0x65, 0x1D)
WARN_TINT = RGBColor(0xFB, 0xEE, 0xE0)
BLUE = RGBColor(0x3A, 0x6E, 0xA5)
BLUE_TINT = RGBColor(0xE7, 0xEE, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDE, 0xE2, 0xD6)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.25)
TAKEAWAY_H = Inches(0.78)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid(); bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or RULE; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, left, top, width, height, color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.12
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color; shp.line.width = Pt(1.0)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic; r.font.name = FONT
    return box


def add_bullets(slide, left, top, width, height, items, size=16, color=INK, space_after=9, bullet_color=ACCENT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.line_spacing = 1.06
        r0 = p.add_run(); r0.text = "▪  "; r0.font.size = Pt(size); r0.font.color.rgb = bullet_color; r0.font.name = FONT
        r1 = p.add_run(); r1.text = item; r1.font.size = Pt(size); r1.font.color.rgb = color; r1.font.name = FONT
    return box


def add_eyebrow(slide, text):
    add_text(slide, MARGIN, Inches(0.42), Inches(12), Inches(0.32), [(text, 13, ACCENT, True, False)])


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    add_text(slide, MARGIN, Inches(0.72), Inches(12.2), Inches(0.62), [(title, 23, INK, True, False)])
    add_rect(slide, MARGIN, Inches(1.38), Inches(1.0), Pt(3), ACCENT)


def add_takeaway(slide, text, tint=ACCENT_TINT, dark=ACCENT_DARK, bar=ACCENT):
    add_rect(slide, 0, TAKEAWAY_TOP, SLIDE_W, TAKEAWAY_H, tint)
    add_rect(slide, 0, TAKEAWAY_TOP, Inches(0.12), TAKEAWAY_H, bar)
    add_text(slide, Inches(0.45), TAKEAWAY_TOP, SLIDE_W - Inches(0.9), TAKEAWAY_H,
              [("TAKEAWAY   ", 12, dark, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.55), TAKEAWAY_TOP, SLIDE_W - Inches(2.0), TAKEAWAY_H,
              [(text, 14.5, dark, True, False)], anchor=MSO_ANCHOR.MIDDLE)


def page_num(slide, n):
    add_text(slide, SLIDE_W - Inches(0.7), Inches(0.42), Inches(0.4), Inches(0.3),
              [(str(n), 11, FAINT, False, False)], align=PP_ALIGN.RIGHT)


def content_box(has_takeaway=True):
    bottom = (TAKEAWAY_TOP - Inches(0.15)) if has_takeaway else (SLIDE_H - Inches(0.3))
    return CONTENT_TOP, bottom - CONTENT_TOP


def styled_table(slide, left, top, width, height, header, rows, col_weights, header_size=12, body_size=11.5, highlight_rows=()):
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, len(header), left, top, width, height).table
    for c, w in enumerate(col_weights):
        tbl.columns[c].width = Emu(int(width * w))
    for ci, text in enumerate(header):
        cell = tbl.cell(0, ci); cell.text = text
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(header_size); p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = FONT
        p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    for ri, row in enumerate(rows, start=1):
        hl = (ri - 1) in highlight_rows
        for ci, text in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = text
            cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_TINT if hl else (WHITE if ri % 2 else RGBColor(0xF5, 0xF6, 0xF1))
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(body_size); p.runs[0].font.bold = (ci == 0 or hl)
            p.runs[0].font.color.rgb = INK; p.runs[0].font.name = FONT
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    return tbl


def flow_box(slide, left, top, width, height, text, fill, text_color=WHITE, size=11):
    shp = add_rounded(slide, left, top, width, height, fill)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(4); tf.margin_right = Pt(4); tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = text_color; r.font.bold = (i == 0); r.font.name = FONT
    return shp


# ============================================================ SLIDE 1: TITLE
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), ACCENT)
add_text(s, MARGIN, Inches(0.26), Inches(12.2), Inches(0.65),
         [("Global ERA5 Data Preparation for Chronos-2 LAI Forecasting", 24, WHITE, True, False)])
page_num(s, 1)
top, h = content_box()
add_text(s, MARGIN, top + Inches(0.1), Inches(11.8), Inches(0.4), [("WEEKLY RESEARCH UPDATE", 13, ACCENT, True, False)])
add_bullets(s, MARGIN, top + Inches(0.65), Inches(11.8), Inches(3.6), [
    "This week's focus: preparing ERA5 global meteorological data to extend our Chronos-2 LAI "
    "forecasting framework beyond the continental U.S.",
    "Covered here: why ERA5, how it maps onto our existing 7-variable pipeline, the global non-U.S. "
    "pixel selection, current download progress, and next steps.",
], size=17, space_after=14)
add_takeaway(s, "Pipeline and experiment design are essentially ready; ERA5 data download is the one remaining dependency before running the global experiment.")


# ============================================================ SLIDE 2: MOTIVATION
s = add_slide(); set_bg(s)
add_title(s, "Extending Beyond gridMET's U.S.-Only Coverage", eyebrow="MOTIVATION")
page_num(s, 2)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(11.8), h, [
    "Every Chronos-2 LAI experiment so far has used gridMET, a daily ~4km meteorological product "
    "that only covers the continental United States.",
    "Professor Wang suggested using ERA5 instead, so we can test whether zero-shot Chronos-2 "
    "generalizes to vegetation and climate regimes well outside the U.S.",
    "ERA5 (ECMWF's 5th-generation reanalysis) provides the same kind of variables at a consistent "
    "0.25° resolution anywhere on Earth, which lets us reuse one common data source globally.",
    "Research question: how well does zero-shot Chronos-2 forecast LAI across diverse global "
    "vegetation and climate regimes when ERA5 supplies the meteorological input?",
], size=17, space_after=16)
add_takeaway(s, "ERA5 is the key that unlocks testing our framework anywhere in the world, not just the U.S.")


# ============================================================ SLIDE 3: VARIABLE MAPPING
s = add_slide(); set_bg(s)
add_title(s, "Mapping ERA5 to Our Existing 7 Climate Variables", eyebrow="DATA SOURCE & VARIABLE MAPPING")
page_num(s, 3)
top, h = content_box()
add_text(s, MARGIN, top, Inches(11.8), Inches(0.6),
         [("Product: ", 14, MUTED, True, False), ("derived-era5-single-levels-daily-statistics", 14, INK, True, False),
          (" (Copernicus Climate Data Store) — global 0.25° reanalysis with server-side daily "
           "aggregation, so we never download raw hourly data.", 14, MUTED, False, False)])
tbl_top = top + Inches(0.75)
header = ["Pipeline variable", "ERA5 source", "How it's obtained"]
rows = [
    ["tmmx (max temp)", "2m_temperature", "direct (daily maximum)"],
    ["tmmn (min temp)", "2m_temperature", "direct (daily minimum)"],
    ["pr (precipitation)", "total_precipitation", "direct (daily sum, m → mm)"],
    ["srad (solar radiation)", "surface_solar_radiation_downwards", "direct (daily mean, J/m² → W/m²)"],
    ["vpd (vapor pressure deficit)", "2m_temperature + 2m_dewpoint_temperature", "derived, FAO-56 Penman-Monteith formula"],
    ["sph (specific humidity)", "2m_dewpoint_temperature + surface_pressure", "derived"],
    ["vs (wind speed)", "10m u- and v-wind components", "magnitude of daily-mean vector"],
]
styled_table(s, MARGIN, tbl_top, SLIDE_W - 2 * MARGIN, Inches(3.7), header, rows, col_weights=[0.30, 0.42, 0.28], header_size=13, body_size=12.5)
add_takeaway(s, "Every mapping and unit conversion was verified against real downloaded values first, not assumed from documentation.")


# ============================================================ SLIDE 4: PIPELINE INTEGRATION
s = add_slide(); set_bg(s)
add_title(s, "Plugging ERA5 into the Existing Pipeline, Unchanged", eyebrow="DATA PROCESSING")
page_num(s, 4)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(11.8), h, [
    "ERA5 daily values are aligned to each LAI 8-day composite window using the exact same "
    "windowed-mean convention our gridMET pipeline already uses — not a naive calendar-week average.",
    "The core Chronos-2 input-building and zero-shot prediction code required zero changes — only a "
    "new ERA5 data-preparation adapter was added.",
    "This keeps ERA5-based results directly comparable to our existing gridMET results, and keeps the "
    "design open to future meteorological sources (e.g., Earth System Model output).",
], size=17, space_after=16)
add_takeaway(s, "ERA5 is a drop-in replacement for gridMET at the input layer — the model and evaluation code are untouched.")


# ============================================================ SLIDE 5: GLOBAL PIXEL SELECTION
s = add_slide(); set_bg(s)
add_title(s, "45 Representative Non-U.S. Pixels Selected", eyebrow="GLOBAL EXPERIMENT SETUP")
page_num(s, 5)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(11.8), Inches(2.3), [
    "Selected from the global ESA CCI Plant Functional Type (PFT) product (300m, true global "
    "coverage) — not the CONUS-only copy used in prior experiments.",
    "Chosen by farthest-point sampling in PFT-composition space, combined with a minimum 800km "
    "geographic separation — not hand-picked or random similar sites.",
    "The continental U.S., Alaska, and Hawaii are explicitly excluded.",
], size=16, space_after=12)
stat_top = top + Inches(2.55)
stats = [("45", "non-U.S. pixels"), ("10", "PFT classes represented"), ("13", "world regions covered"), ("0.44–1.00", "PFT purity range")]
box_w = Inches(2.75); gap = Inches(0.25)
for i, (num, label) in enumerate(stats):
    left = MARGIN + i * (box_w + gap)
    add_rounded(s, left, stat_top, box_w, Inches(1.15), ACCENT_TINT)
    add_text(s, left, stat_top + Inches(0.12), box_w, Inches(0.55), [(num, 26, ACCENT_DARK, True, False)], align=PP_ALIGN.CENTER)
    add_text(s, left, stat_top + Inches(0.68), box_w, Inches(0.4), [(label, 11.5, MUTED, False, False)], align=PP_ALIGN.CENTER)
add_text(s, MARGIN, stat_top + Inches(1.4), Inches(11.8), Inches(0.5),
         [("Top regions: Siberia/Boreal Eurasia (9), Canada (5), East Asia (4), Amazon (3), Western Europe (3), "
           "East Africa/Sahel (2), Central/Southern Africa (2)", 13, MUTED, False, True)])
add_takeaway(s, "Coverage spans forests, grasslands, and shrublands across every inhabited continent.")


# ============================================================ SLIDE 6: DOWNLOAD PROGRESS & CONSTRAINTS
s = add_slide(); set_bg(s)
add_title(s, "ERA5 Download In Progress — Working Through CDS Constraints", eyebrow="CURRENT STATUS")
page_num(s, 6)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(11.8), h, [
    "Currently validating the ERA5 pipeline on an existing CONUS pixel before scaling to the 45 "
    "global pixels, to catch any unit or alignment errors early.",
    "The Copernicus Climate Data Store (CDS) imposes two hard, undocumented limits found through "
    "direct testing: only one request in flight per account at a time, and a tight per-request size "
    "cap (effectively one variable, one year, per request).",
    "Per-request queue wait time varies widely (minutes to tens of minutes) and is controlled "
    "entirely by CDS, not by us.",
    "This week also included a multi-hour, account-wide CDS service outage — confirmed by testing "
    "unrelated requests and cross-checking the ECMWF community forum — independent of our pipeline.",
    "The download is fully resumable: every (variable, year, location) combination is cached "
    "individually, so no progress is ever lost across interruptions.",
], size=15.5, space_after=11)
add_takeaway(s, "The bottleneck is CDS's own queue and rate limits, not our pipeline — the resumable design keeps making steady progress regardless.")


# ============================================================ SLIDE 7: STATUS SUMMARY
s = add_slide(); set_bg(s)
add_title(s, "What's Done, What's Still Running", eyebrow="STATUS SUMMARY")
page_num(s, 7)
top, h = content_box()
col_w = Emu(int((SLIDE_W - 2 * MARGIN - Inches(0.4)) / 2))
left_col = MARGIN
right_col = MARGIN + col_w + Inches(0.4)

add_rounded(s, left_col, top, col_w, Inches(0.5), ACCENT)
add_text(s, left_col, top, col_w, Inches(0.5), [("COMPLETED", 15, WHITE, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, left_col, top + Inches(0.65), col_w, Inches(3.8), [
    "ERA5 variable mapping and unit conversions (verified empirically)",
    "ERA5-to-LAI temporal alignment logic",
    "Chronos-2 integration (reuses the existing pipeline unchanged)",
    "Global non-U.S. pixel selection (45 pixels)",
], size=15, space_after=12, bullet_color=ACCENT)

add_rounded(s, right_col, top, col_w, Inches(0.5), WARN)
add_text(s, right_col, top, col_w, Inches(0.5), [("STILL IN PROGRESS", 15, WHITE, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_bullets(s, right_col, top + Inches(0.65), col_w, Inches(3.8), [
    "ERA5 download for the pipeline validation pixel",
    "Full 45-pixel global ERA5 download (not yet started)",
    "Global zero-shot Chronos-2 evaluation (blocked on the download above)",
], size=15, space_after=12, bullet_color=WARN)
add_takeaway(s, "Data preparation is essentially complete — the remaining work is waiting on CDS downloads, not new design or engineering.")


# ============================================================ SLIDE 8: NEXT STEPS
s = add_slide(); set_bg(s)
add_title(s, "From Data Preparation to Global Experiments", eyebrow="NEXT STEPS")
page_num(s, 8)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(11.8), h, [
    "Finish the ERA5 download for the validation pixel and compare zero-shot Chronos-2 results "
    "against the existing gridMET baseline.",
    "Once validated, launch the ERA5 download for the full 45-pixel global set.",
    "Run zero-shot Chronos-2 across all 45 pixels and analyze results by vegetation type, climate "
    "regime, and region.",
    "Compare global (ERA5) generalization performance against our existing U.S. (gridMET) results.",
], size=17, space_after=16)
add_takeaway(s, "Next week's goal: move from data preparation into running and analyzing the global Chronos-2 experiment itself.")

out_path = f"{ROOT}/ERA5_Global_Chronos2_Progress_Update.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
