# -*- coding: utf-8 -*-
"""Builds the research-group-update deck. All numbers/claims are pulled
directly from LOYO_CV_FINDINGS.md, PREDICTOR_ABLATION_REPORT.md, and
FINETUNE_IMPROVEMENT_REPORT.md (re-read immediately before writing this
script) - nothing here is invented."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import copy

from pathlib import Path
ROOT = str(Path(__file__).resolve().parent.parent.parent)
CHRONOS2 = f"{ROOT}/Chronos2-vegetation-forecasting/outputs"
AELSTM = f"{ROOT}/AELSTM/outputs"

# ---- palette / type ----
INK = RGBColor(0x1C, 0x21, 0x1C)
MUTED = RGBColor(0x5B, 0x63, 0x58)
FAINT = RGBColor(0x8B, 0x93, 0x85)
ACCENT = RGBColor(0x2F, 0x6F, 0x5E)
ACCENT_DARK = RGBColor(0x1E, 0x4A, 0x3E)
ACCENT_TINT = RGBColor(0xEA, 0xF2, 0xEF)
NEG = RGBColor(0xA3, 0x40, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDD, 0xE2, 0xDC)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.25)
TAKEAWAY_H = Inches(0.78)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = RULE
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, wrap=True):
    """runs: list of (text, size, color, bold, italic) tuples, one per paragraph."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
    return box


def add_bullets(slide, left, top, width, height, items, size=17, color=INK, space_after=10,
                 bullet_color=ACCENT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        # bullet glyph as separate colored run
        r0 = p.add_run()
        r0.text = "▪  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        r0.font.name = FONT
        r1 = p.add_run()
        r1.text = item
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.name = FONT
    return box


def add_eyebrow(slide, part_label, kicker=None):
    txt = part_label if not kicker else f"{part_label}   •   {kicker}"
    add_text(slide, MARGIN, Inches(0.42), Inches(11), Inches(0.32),
              [(txt, 13, ACCENT, True, False)])


def add_title(slide, title, part_label=None, kicker=None):
    if part_label:
        add_eyebrow(slide, part_label, kicker)
    add_text(slide, MARGIN, Inches(0.72), Inches(12.2), Inches(0.62),
              [(title, 27, INK, True, False)])
    add_rect(slide, MARGIN, Inches(1.38), Inches(1.0), Pt(3), ACCENT)


def add_takeaway(slide, text):
    add_rect(slide, 0, TAKEAWAY_TOP, SLIDE_W, TAKEAWAY_H, ACCENT_TINT)
    add_rect(slide, 0, TAKEAWAY_TOP, Inches(0.12), TAKEAWAY_H, ACCENT)
    add_text(slide, Inches(0.45), TAKEAWAY_TOP, SLIDE_W - Inches(0.9), TAKEAWAY_H,
              [("TAKEAWAY   ", 12, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.55), TAKEAWAY_TOP, SLIDE_W - Inches(2.0), TAKEAWAY_H,
              [(text, 16.5, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)


def page_num(slide, n):
    add_text(slide, SLIDE_W - Inches(0.7), Inches(0.42), Inches(0.4), Inches(0.3),
              [(str(n), 11, FAINT, False, False)], align=PP_ALIGN.RIGHT)


def add_picture_fit(slide, path, box_left, box_top, box_w, box_h, align="center"):
    im = Image.open(path)
    ar = im.size[0] / im.size[1]
    box_ar = box_w / box_h
    if ar > box_ar:
        w = box_w
        h = int(w / ar)
    else:
        h = box_h
        w = int(h * ar)
    if align == "center":
        left = box_left + (box_w - w) // 2
    else:
        left = box_left
    top = box_top + (box_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    return left, top, w, h


def content_box(has_takeaway=True):
    bottom = (TAKEAWAY_TOP - Inches(0.15)) if has_takeaway else (SLIDE_H - Inches(0.3))
    return CONTENT_TOP, bottom - CONTENT_TOP


# ============================================================ SLIDE 1: TITLE
s = add_slide()
set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(7.5), WHITE)
add_rect(s, 0, Inches(4.55), SLIDE_W, Inches(2.95), ACCENT)
add_text(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.4),
          [("RESEARCH GROUP UPDATE", 15, ACCENT, True, False)])
add_text(s, Inches(0.9), Inches(1.98), Inches(11.6), Inches(1.9),
          [("Vegetation LAI Forecasting:", 38, INK, True, False),
           ("Robustness, Predictors, and Fine-Tuning", 38, INK, True, False)], line_spacing=1.05)
add_rect(s, Inches(0.9), Inches(3.75), Inches(1.3), Pt(4), ACCENT)
add_text(s, Inches(0.9), Inches(3.95), Inches(11), Inches(0.5),
          [("AELSTM + Chronos-2 vegetation forecasting project", 17, MUTED, False, True)])
add_text(s, Inches(0.9), Inches(5.05), Inches(11.5), Inches(2.0),
          [("Rolling-Origin LOYO Cross-Validation", 19, WHITE, True, False),
           ("Predictor Sensitivity / Ablation Study", 19, WHITE, True, False),
           ("Improved Chronos-2 LoRA Fine-Tuning", 19, WHITE, True, False)], line_spacing=1.35)
add_text(s, Inches(0.9), Inches(7.0), Inches(11), Inches(0.35),
          [("August 2026", 12, RGBColor(0xD8, 0xE6, 0xE1), False, False)])

# ============================================================ SLIDE 2: RECAP
s = add_slide(); set_bg(s)
add_title(s, "Recap: Pixels and Baseline Setup")
page_num(s, 2)
top, h = content_box()
img_w, img_h = Inches(5.6), Inches(4.4)
add_picture_fit(s, f"{AELSTM}/pixel_selection/pixel_locations_map.png", MARGIN, top, img_w, img_h, align="left")

tbl_left = MARGIN + img_w + Inches(0.4)
tbl_w = SLIDE_W - tbl_left - MARGIN
rows_data = [
    ("Pixel", "Location", "Land cover"),
    ("low_amplitude", "37.53°N, 117.56°W", "Natural grassland"),
    ("high_amplitude_\ndeciduous", "36.23°N, 84.47°W", "Broadleaf deciduous forest"),
    ("evergreen", "30.53°N, 82.43°W", "Needleleaf evergreen forest"),
]
gtbl = s.shapes.add_table(4, 3, tbl_left, top, tbl_w, Inches(1.7)).table
gtbl.columns[0].width = Inches(1.9)
gtbl.columns[1].width = Inches(2.0)
gtbl.columns[2].width = tbl_w - Inches(3.9)
for ci in range(3):
    cell = gtbl.cell(0, ci)
    cell.text = rows_data[0][ci]
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(13); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = FONT
    cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)
for ri in range(1, 4):
    for ci in range(3):
        cell = gtbl.cell(ri, ci)
        cell.text = rows_data[ri][ci]
        cell.fill.solid(); cell.fill.fore_color.rgb = (ACCENT_TINT if ri % 2 == 0 else WHITE)
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(12.5); p.runs[0].font.bold = (ci == 0)
        p.runs[0].font.color.rgb = INK; p.runs[0].font.name = FONT
        cell.margin_top = Pt(4); cell.margin_bottom = Pt(4)

add_bullets(s, tbl_left, top + Inches(2.0), tbl_w, Inches(2.6), [
    "8 AELSTM-family models: AELSTM, BiLSTM, LSTM, GRU, RNN, CNN, RF, SVM",
    "Chronos-2 (zero-shot and LoRA fine-tuned), using historical LAI + historical climate + known future climate",
    "7 gridMET predictors: tmmx, tmmn, pr, srad, vpd, sph, vs",
    "Baseline protocol: train 2000–2021, test 2022, scored against raw observed LAI",
], size=15.5)
add_takeaway(s, "Same 3 pixels, same 8+2 methods, and raw-observation scoring underlie all 3 studies below.")

# ============================================================ PART 1 SLIDES
PART1 = "PART 1 — ROLLING-ORIGIN LOYO-CV"

# ---- Slide 3: motivation
s = add_slide(); set_bg(s)
add_title(s, "Motivation: Is One Test Year Enough?", PART1)
page_num(s, 3)
top, h = content_box()
add_bullets(s, MARGIN, top + Inches(0.3), Inches(12.2), h, [
    "The baseline protocol trains on 2000–2021 and tests on exactly one year: 2022.",
    "A single test year cannot separate “good model” from “model that got a lucky (or unlucky) year.”",
    "2022 turned out not to be an anomalous year at any of the 3 pixels — so single-year evaluation gives no evidence about how models behave when a real anomaly (drought, noisy year) occurs.",
], size=19, space_after=22)
add_takeaway(s, "We need evidence across many held-out years, not confidence built on one.")

# ---- Slide 4: LOYO-CV design
s = add_slide(); set_bg(s)
add_title(s, "Rolling-Origin LOYO-CV: Design", PART1)
page_num(s, 4)
top, h = content_box()
add_bullets(s, MARGIN, top, Inches(12.2), Inches(1.5), [
    "Fixed 12-year training window immediately before each held-out year (not classic all-other-years LOYO, and not an expanding window).",
    "11 folds: held-out years 2012–2022. Same 10 methods, same 3 pixels, scored against raw observed LAI.",
], size=17.5, space_after=12)

# schematic: 4 example fold bars
sch_top = top + Inches(1.75)
labels = [("2000–2011", "2012"), ("2001–2012", "2013"), ("...", "..."), ("2010–2021", "2022")]
bar_w = Inches(2.55)
gap = Inches(0.35)
start_x = MARGIN + Inches(0.3)
def emu(v):
    return Emu(int(v))

for i, (train_lbl, test_lbl) in enumerate(labels):
    x = emu(start_x + i * (bar_w + gap))
    add_rect(s, x, sch_top, bar_w, Inches(0.55), ACCENT_TINT)
    add_text(s, x, sch_top, bar_w, Inches(0.55), [(f"train {train_lbl}", 12.5, ACCENT_DARK, True, False)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, emu(x + bar_w * 0.32), sch_top + Inches(0.62), emu(bar_w * 0.36), Inches(0.28))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED; arrow.line.fill.background(); arrow.shadow.inherit = False
    add_rect(s, emu(x + bar_w * 0.15), sch_top + Inches(1.05), emu(bar_w * 0.7), Inches(0.5), ACCENT)
    add_text(s, emu(x + bar_w * 0.15), sch_top + Inches(1.05), emu(bar_w * 0.7), Inches(0.5),
              [(f"test {test_lbl}", 12.5, WHITE, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, sch_top + Inches(1.7), bar_w, Inches(0.3), [(f"fold {i+1 if i<3 else 11}", 11, FAINT, False, True)],
              align=PP_ALIGN.CENTER)

add_takeaway(s, "A fixed window isolates genuine year-to-year difficulty from training-set-size effects.")

# ---- Slide 5: main LOYO results
s = add_slide(); set_bg(s)
add_title(s, "Model Performance: Stable at Some Pixels, Volatile at Others", PART1)
page_num(s, 5)
top, h = content_box()
add_picture_fit(s, f"{CHRONOS2}/loyo_cv/comparison/loyo_r2_distributions.png",
                  MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(3.15), align="center")
add_bullets(s, MARGIN, top + Inches(3.35), SLIDE_W - 2 * MARGIN, Inches(1.3), [
    "high_amplitude_deciduous: all 10 methods score R² 0.89–0.99 across 11 years — stable regardless of method.",
    "low_amplitude and evergreen: much wider spread, and every method has one severe outlier fold.",
    "At evergreen, every method’s mean R² is near zero or negative — but every median is a healthy 0.48–0.68.",
], size=16, space_after=8)
add_takeaway(s, "Report median alongside mean — a single anomalous fold can dominate the average.")

# ---- Slide 6: ranking consistency
s = add_slide(); set_bg(s)
add_title(s, "Which Models Are Consistently Strong?", PART1)
page_num(s, 6)
top, h = content_box()
add_picture_fit(s, f"{CHRONOS2}/loyo_cv/comparison/loyo_ranking_frequency.png",
                  MARGIN, top, Inches(7.6), h - Inches(0.1), align="left")
add_bullets(s, MARGIN + Inches(7.9), top + Inches(0.3), Inches(4.5), h, [
    "RF (mean rank 3.39) and Chronos-2 zero-shot (3.70) are top-3 on ~6 of every 10 folds, and rarely worst.",
    "LoRA fine-tuning is the single most volatile method — frequent top-3 AND frequent bottom-3 finishes.",
    "AELSTM (bottom-3 on 36% of folds) and CNN (58%) are the least reliable methods overall.",
], size=15.5, space_after=14)
add_takeaway(s, "Random Forest and Chronos-2 zero-shot are the most dependable models across all 33 folds.")

# ---- Slide 7: outlier analysis
s = add_slide(); set_bg(s)
add_title(s, "Two Outlier Folds, Two Different Failure Modes", PART1)
page_num(s, 7)
top, h = content_box()
add_picture_fit(s, f"{CHRONOS2}/loyo_cv/comparison/loyo_outlier_fold_investigation.png",
                  MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(3.25), align="center")
add_bullets(s, MARGIN, top + Inches(3.45), SLIDE_W - 2 * MARGIN, Inches(1.2), [
    "evergreen / 2012 (z = −3.0): a real, sustained drought — LAI never reaches its expected seasonal peak for months.",
    "low_amplitude / 2018 (z = −2.9): not a climate anomaly — an erratic, noisy year at an already low-signal pixel.",
    "No year was difficult across all 3 pixels at once — difficulty is pixel-specific, not a shared climate signal.",
], size=15.5, space_after=8)
add_takeaway(s, "A single fixed 2022 test year could never have revealed either failure mode.")

# ---- Slide 8: Part 1 takeaway
s = add_slide(); set_bg(s)
add_title(s, "Part 1 Takeaway: LOYO-CV", PART1)
page_num(s, 8)
top, h = content_box()
add_bullets(s, MARGIN, top + Inches(0.2), Inches(12.2), h, [
    "Rolling-origin LOYO-CV gives a far more robust, realistic evaluation than a single fixed test year.",
    "Model performance varies substantially across pixels and years — there is no universal winner.",
    "Random Forest and Chronos-2 zero-shot are the most consistently strong models across 33 folds.",
    "Chronos-2’s advantage comes from normal years, not special resilience to anomalous ones — it was among the worst methods on the evergreen drought fold.",
], size=18, space_after=18)
add_takeaway(s, "LOYO-CV reveals failure modes that single-year (2022-only) evaluation completely misses.")

# ============================================================ PART 2 SLIDES
PART2 = "PART 2 — PREDICTOR SENSITIVITY"

# ---- Slide 9: design
s = add_slide(); set_bg(s)
add_title(s, "Predictor Sensitivity: Experimental Design", PART2)
page_num(s, 9)
top, h = content_box()
add_bullets(s, MARGIN, top + Inches(0.2), Inches(12.2), Inches(2.6), [
    "Baseline: all 7 predictors (tmmx, tmmn, pr, srad, vpd, sph, vs) — reused from already-verified results, not rerun.",
    "Phase 1 — leave-one-predictor-out: 7 configs, each dropping exactly one predictor.",
    "Phase 2 — 4 grouped ablations, selection rule fixed before Phase 1 ran: top-3 essential only, drop the least-important pair, plus two a priori pairs — {tmmx,tmmn} and {vpd,sph}.",
    "9 methods (8 AELSTM-family + Chronos-2 zero-shot) × 3 pixels; Chronos-2 fine-tuning excluded to keep the search computationally manageable.",
], size=17, space_after=16)
add_takeaway(s, "A small, pre-registered design instead of an intractable 2⁷ = 128-subset search.")

# ---- Slide 10: heatmap
s = add_slide(); set_bg(s)
add_title(s, "Predictor Importance Is Model- and Pixel-Specific", PART2)
page_num(s, 10)
top, h = content_box()
add_picture_fit(s, f"{CHRONOS2}/predictor_ablation/comparison/predictor_model_pixel_heatmap.png",
                  MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(3.15), align="center")
add_bullets(s, MARGIN, top + Inches(3.35), SLIDE_W - 2 * MARGIN, Inches(1.3), [
    "Same predictor removal helps one model and hurts another at the same pixel (e.g. dropping pr helps GRU, hurts RNN/CNN at low_amplitude).",
    "high_amplitude_deciduous is largely insensitive to any single predictor — near-ceiling performance regardless.",
    "low_amplitude and evergreen show large, heterogeneous ΔR² in both directions.",
], size=16, space_after=8)
add_takeaway(s, "No single predictor-importance ranking applies to every model or every pixel.")

# ---- Slide 11: key findings Part 2
s = add_slide(); set_bg(s)
add_title(s, "srad Matters Most; vs and pr Are Redundant", PART2)
page_num(s, 11)
top, h = content_box()
add_picture_fit(s, f"{CHRONOS2}/predictor_ablation/comparison/predictor_reduced_sets.png",
                  MARGIN, top, Inches(7.6), h - Inches(0.1), align="left")
add_bullets(s, MARGIN + Inches(7.9), top + Inches(0.3), Inches(4.5), h, [
    "srad is the only predictor that helps at all 3 pixels — the single most important variable overall.",
    "Temperature (tmmx/tmmn) is critical at low_amplitude but unimportant — even mildly harmful — elsewhere.",
    "Dropping the 2 least-important predictors (vs, pr) improves mean R² at all 3 pixels simultaneously.",
], size=15.5, space_after=14)
add_takeaway(s, "A leaner 5-predictor set (tmmx, tmmn, srad, vpd, sph) is a validated, better default.")

# ============================================================ PART 3 SLIDES
PART3 = "PART 3 — IMPROVED LoRA FINE-TUNING"

# ---- Slide 12: why unsatisfactory
s = add_slide(); set_bg(s)
add_title(s, "The Original Fine-Tuning Had No Safety Net", PART3)
page_num(s, 12)
top, h = content_box()
add_bullets(s, MARGIN, top + Inches(0.3), Inches(12.2), h, [
    "fit() was called without validation_inputs → eval_strategy=“no”, load_best_model_at_end=False.",
    "Trained a fixed 1000 steps and kept whatever weights existed at that instant — no way to detect overfitting.",
    "LoRA hyperparameters (lr=1e-4, r=8, α=16) were taken verbatim from Amazon’s own notebook, never tuned for this dataset.",
], size=19, space_after=22)
add_takeaway(s, "No validation signal → no way to know whether fine-tuning helped or hurt.")

# ---- Slide 13: improved design
s = add_slide(); set_bg(s)
add_title(s, "Improved Design: Validation-Based Checkpoint Selection", PART3)
page_num(s, 13)
top, h = content_box()
add_picture_fit(s, f"{CHRONOS2}/finetuned_lora_improved/comparison/finetune_winner_train_val_curves.png",
                  MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(3.05), align="center")
add_bullets(s, MARGIN, top + Inches(3.25), SLIDE_W - 2 * MARGIN, Inches(1.4), [
    "Chronological validation split from pre-2022 data (2020 & 2021 folds) — 2022 stays completely untouched.",
    "Small search: learning rate ∈ {1e-5, 1e-4} × LoRA rank ∈ {4, 8, 16}, with early stopping; best checkpoint chosen purely by validation loss.",
    "low_amplitude overfits immediately (best step = 100); the other two pixels show genuine interior validation-loss minima.",
], size=15.5, space_after=8)
add_takeaway(s, "Confirmed from source: validation_inputs auto-selects the best checkpoint, not the final training step.")

# ---- Slide 14: comparison
s = add_slide(); set_bg(s)
add_title(s, "Validation-Based Tuning Recovers Most of the Lost Performance", PART3)
page_num(s, 14)
top, h = content_box()

tbl_w2 = Inches(6.0)
rows2 = [
    ("Pixel", "Zero-shot", "Original\nLoRA", "Improved\nLoRA"),
    ("low_amplitude", "0.542", "0.379", "0.500"),
    ("high_amplitude_\ndeciduous", "0.965", "0.952", "0.942"),
    ("evergreen", "0.832", "0.712", "0.829"),
]
t2 = s.shapes.add_table(4, 4, MARGIN, top + Inches(0.1), tbl_w2, Inches(2.1)).table
t2.columns[0].width = Inches(2.2)
remaining_w = tbl_w2 - Inches(2.2)  # already in EMU - do not wrap in Inches() again
for c in range(1, 4):
    t2.columns[c].width = Emu(int(remaining_w / 3))
for ci in range(4):
    cell = t2.cell(0, ci)
    cell.text = rows2[0][ci]
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    p = cell.text_frame.paragraphs[0]; p.runs[0].font.size = Pt(13); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = WHITE; p.runs[0].font.name = FONT
    p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
for ri in range(1, 4):
    for ci in range(4):
        cell = t2.cell(ri, ci)
        cell.text = rows2[ri][ci]
        cell.fill.solid(); cell.fill.fore_color.rgb = (ACCENT_TINT if ri % 2 == 0 else WHITE)
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size = Pt(14); p.runs[0].font.bold = (ci == 0)
        p.runs[0].font.color.rgb = INK; p.runs[0].font.name = FONT
        p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT

add_bullets(s, MARGIN, top + Inches(2.55), tbl_w2, Inches(1.6), [
    "Gap to zero-shot shrinks from −0.163 to −0.042 at low_amplitude, and from −0.120 to −0.003 at evergreen.",
    "Improved LoRA still does not outright beat zero-shot on any of the 3 pixels.",
], size=15, space_after=9)

img_left = MARGIN + tbl_w2 + Inches(0.35)
img_w2 = SLIDE_W - img_left - MARGIN
add_picture_fit(s, f"{CHRONOS2}/finetuned_lora_improved/comparison/three_way_comparison_evergreen.png",
                  img_left, top + Inches(0.1), img_w2, Inches(3.9), align="left")
add_takeaway(s, "Most of the original degradation was a training-setup artifact, not fundamental to LoRA itself.")

# ============================================================ SLIDE 15: OVERALL CONCLUSIONS
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.15), ACCENT)
add_text(s, MARGIN, Inches(0.32), Inches(11), Inches(0.6),
          [("Overall Conclusions & Next Steps", 26, WHITE, True, False)])
page_num(s, 15)

col_w = Emu(int((SLIDE_W - 2 * MARGIN - Inches(0.5)) / 2))
top2 = Inches(1.55)
add_text(s, MARGIN, top2, col_w, Inches(0.4), [("What We Learned", 17, ACCENT_DARK, True, False)])
add_bullets(s, MARGIN, top2 + Inches(0.5), col_w, Inches(5.0), [
    "Rolling-origin LOYO-CV is essential — it reveals failure modes a single 2022 test year cannot.",
    "Model performance and predictor importance both vary substantially across pixels and years.",
    "srad is the most consistently important predictor; vs and pr are measurably redundant.",
    "Random Forest and Chronos-2 zero-shot are the most consistently strong models across LOYO folds — Chronos-2 zero-shot remains very competitive.",
    "The original LoRA degradation was largely caused by the fine-tuning setup, not LoRA itself; validation-based checkpoint selection recovers most of it — but improved LoRA still doesn’t consistently beat zero-shot.",
], size=14.5, space_after=13)

col2_left = MARGIN + col_w + Inches(0.5)
add_text(s, col2_left, top2, col_w, Inches(0.4), [("Next Steps", 17, ACCENT_DARK, True, False)])
add_bullets(s, col2_left, top2 + Inches(0.5), col_w, Inches(5.0), [
    "Extend LOYO-CV and the predictor ablation to the remaining 3 selected pixels (cropland, grassland/shrubland, high-interannual-variability).",
    "Adopt the reduced 5-predictor set (tmmx, tmmn, srad, vpd, sph) as the new default and re-validate across more pixels.",
    "Explore validation-informed fine-tuning further — more validation folds, a wider search — to test whether LoRA can be made to consistently beat zero-shot.",
    "Investigate why Chronos-2 (zero-shot and fine-tuned) underperforms specifically on sustained climate anomalies like the evergreen drought fold.",
], size=14.5, space_after=13, bullet_color=ACCENT)

out_path = f"{ROOT}/Chronos2-vegetation-forecasting/Vegetation_Forecasting_Research_Update.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
