# -*- coding: utf-8 -*-
"""Builds a short 5-slide summary deck for Prof. Wang's leakage-diagnostic
follow-up question. Reuses the exact visual identity (palette, fonts, layout
helpers) from build_research_update_deck.py. All numbers/figures are pulled
directly from AELSTM/LEAKAGE_DIAGNOSTIC_REPORT.md,
AELSTM/outputs/leakage_diagnostic_2012/evergreen/,
Chronos2-vegetation-forecasting/LEAKAGE_DIAGNOSTIC_REPORT_CHRONOS2.md, and
Chronos2-vegetation-forecasting/outputs/leakage_diagnostic_2012/evergreen/ -
nothing here is invented."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = str(Path(__file__).resolve().parent.parent.parent)
CHRONOS2 = f"{ROOT}/Chronos2-vegetation-forecasting/outputs"
AELSTM = f"{ROOT}/AELSTM/outputs"

# ---- palette / type (matches build_research_update_deck.py) ----
INK = RGBColor(0x1C, 0x21, 0x1C)
MUTED = RGBColor(0x5B, 0x63, 0x58)
FAINT = RGBColor(0x8B, 0x93, 0x85)
ACCENT = RGBColor(0x2F, 0x6F, 0x5E)
ACCENT_DARK = RGBColor(0x1E, 0x4A, 0x3E)
ACCENT_TINT = RGBColor(0xEA, 0xF2, 0xEF)
NEG = RGBColor(0xA3, 0x40, 0x2F)
NEG_TINT = RGBColor(0xF7, 0xE7, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xDD, 0xE2, 0xDC)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.25)
TAKEAWAY_H = Inches(0.78)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def emu(v):
    return Emu(int(v))


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = RULE
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, size, color, bold, italic) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
    return box


def add_bullets(slide, left, top, width, height, items, size=17, color=INK, space_after=10,
                 bullet_color=ACCENT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.08
        r0 = p.add_run()
        r0.text = "▪  "
        r0.font.size = Pt(size)
        r0.font.color.rgb = bullet_color
        r0.font.name = FONT
        r1 = p.add_run()
        r1.text = item
        r1.font.size = Pt(size)
        r1.font.color.rgb = color
        r1.font.name = FONT
    return box


def add_eyebrow(slide, text):
    add_text(slide, MARGIN, Inches(0.42), Inches(12), Inches(0.32),
              [(text, 13, ACCENT, True, False)])


def add_title(slide, title, eyebrow=None):
    if eyebrow:
        add_eyebrow(slide, eyebrow)
    add_text(slide, MARGIN, Inches(0.72), Inches(12.2), Inches(0.62),
              [(title, 27, INK, True, False)])
    add_rect(slide, MARGIN, Inches(1.38), Inches(1.0), Pt(3), ACCENT)


def add_takeaway(slide, text):
    add_rect(slide, 0, TAKEAWAY_TOP, SLIDE_W, TAKEAWAY_H, ACCENT_TINT)
    add_rect(slide, 0, TAKEAWAY_TOP, Inches(0.12), TAKEAWAY_H, ACCENT)
    add_text(slide, Inches(0.45), TAKEAWAY_TOP, SLIDE_W - Inches(0.9), TAKEAWAY_H,
              [("TAKEAWAY   ", 12, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.55), TAKEAWAY_TOP, SLIDE_W - Inches(2.0), TAKEAWAY_H,
              [(text, 16.5, ACCENT_DARK, True, False)], anchor=MSO_ANCHOR.MIDDLE)


def page_num(slide, n):
    add_text(slide, SLIDE_W - Inches(0.7), Inches(0.42), Inches(0.4), Inches(0.3),
              [(str(n), 11, FAINT, False, False)], align=PP_ALIGN.RIGHT)


def add_picture_fit(slide, path, box_left, box_top, box_w, box_h, align="center"):
    im = Image.open(path)
    ar = im.size[0] / im.size[1]
    box_ar = box_w / box_h
    if ar > box_ar:
        w = box_w
        h = int(w / ar)
    else:
        h = box_h
        w = int(h * ar)
    left = box_left + (box_w - w) // 2 if align == "center" else box_left
    top = box_top + (box_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    return left, top, w, h


def content_box(has_takeaway=True):
    bottom = (TAKEAWAY_TOP - Inches(0.15)) if has_takeaway else (SLIDE_H - Inches(0.3))
    return CONTENT_TOP, bottom - CONTENT_TOP


# ============================================================ SLIDE 1: RESEARCH QUESTION
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(7.5), WHITE)
add_rect(s, 0, Inches(0), SLIDE_W, Inches(1.15), ACCENT)
add_text(s, MARGIN, Inches(0.32), Inches(12), Inches(0.6),
          [("Leakage Diagnostic — Evergreen / 2012", 26, WHITE, True, False)])
page_num(s, 1)

top, h = content_box()
add_text(s, MARGIN, top + Inches(0.1), Inches(12.2), Inches(0.4),
          [("PROF. WANG'S QUESTION", 13, ACCENT, True, False)])
add_rect(s, MARGIN, top + Inches(0.55), Inches(0.1), Inches(1.3), ACCENT)
add_text(s, MARGIN + Inches(0.35), top + Inches(0.5), Inches(11.8), Inches(1.4),
          [("“What happens if you predict the 2012 LAI using a model that has "
            "already seen 2012 during training?”", 22, INK, True, True)], line_spacing=1.15)

add_bullets(s, MARGIN, top + Inches(2.35), Inches(12.2), Inches(2.2), [
    "In LOYO-CV, every one of the 10 methods tested collapses to strongly negative R² "
    "(−4.5 to −8.5) on the evergreen pixel's 2012 fold — a real, severe drought year.",
    "Purpose of this diagnostic: find out whether that failure is mainly caused by 2012 being "
    "an unseen, out-of-distribution year, or by a limitation of the models themselves.",
], size=18, space_after=18)
add_takeaway(s, "We test this directly: let each model see 2012 during fitting, and see how much recovers.")

# ============================================================ SLIDE 2: EXPERIMENTAL DESIGN
s = add_slide(); set_bg(s)
add_title(s, "Experimental Design")
page_num(s, 2)
top, h = content_box()

col_w = emu((SLIDE_W - 2 * MARGIN - Inches(0.6)) / 2)
col2_left = emu(MARGIN + col_w + Inches(0.6))

# Left: ORIGINAL
add_rect(s, MARGIN, top, col_w, Inches(0.5), ACCENT)
add_text(s, MARGIN, top, col_w, Inches(0.5), [("ORIGINAL — valid evaluation", 15, WHITE, True, False)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

sch_top = top + Inches(0.85)
bar_h = Inches(0.55)
add_rect(s, MARGIN, sch_top, col_w, bar_h, ACCENT_TINT)
add_text(s, MARGIN, sch_top, col_w, bar_h, [("Train: 2000–2011", 15, ACCENT_DARK, True, False)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, emu(MARGIN + col_w * 0.42), sch_top + bar_h + Inches(0.06),
                            emu(col_w * 0.16), Inches(0.3))
arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED; arrow.line.fill.background(); arrow.shadow.inherit = False
add_rect(s, MARGIN, sch_top + bar_h + Inches(0.45), col_w, bar_h, ACCENT)
add_text(s, MARGIN, sch_top + bar_h + Inches(0.45), col_w, bar_h,
          [("Test: 2012 (unseen)", 15, WHITE, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, MARGIN, sch_top + 2 * bar_h + Inches(0.95), col_w, Inches(0.9),
          [("2012 never participates in fitting — this is the standard LOYO-CV protocol.",
            13.5, MUTED, False, True)], line_spacing=1.2)

# Right: DIAGNOSTIC (leakage)
add_rect(s, col2_left, top, col_w, Inches(0.5), NEG)
add_text(s, col2_left, top, col_w, Inches(0.5), [("DIAGNOSTIC — data leakage", 15, WHITE, True, False)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, col2_left, sch_top, col_w, bar_h, NEG_TINT)
add_text(s, col2_left, sch_top, col_w, bar_h, [("Train: includes 2012", 15, NEG, True, False)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, emu(col2_left + col_w * 0.42), sch_top + bar_h + Inches(0.06),
                             emu(col_w * 0.16), Inches(0.3))
arrow2.fill.solid(); arrow2.fill.fore_color.rgb = MUTED; arrow2.line.fill.background(); arrow2.shadow.inherit = False
add_rect(s, col2_left, sch_top + bar_h + Inches(0.45), col_w, bar_h, NEG)
add_text(s, col2_left, sch_top + bar_h + Inches(0.45), col_w, bar_h,
          [("Evaluate: 2012 (seen)", 15, WHITE, True, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

badge_top = sch_top + 2 * bar_h + Inches(0.95)
add_rect(s, col2_left, badge_top, col_w, Inches(0.55), NEG)
add_text(s, col2_left, badge_top, col_w, Inches(0.55),
          [("⚠  NOT VALID EVALUATION — DELIBERATE DATA LEAKAGE", 13, WHITE, True, False)],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_bullets(s, MARGIN, badge_top + Inches(0.85), SLIDE_W - 2 * MARGIN, Inches(1.1), [
    "8 AELSTM-family models (full refit) + Chronos-2 via LoRA fine-tuning — zero-shot has no "
    "trainable weights, so it cannot “learn” 2012 and is excluded from this condition.",
    "The evaluation forecast call is identical in both conditions — only the fitting data differs.",
], size=14.5, space_after=8)
add_takeaway(s, "The only variable that changes is whether 2012 participated in fitting.")

# ============================================================ SLIDE 3: AELSTM-FAMILY RESULTS
s = add_slide(); set_bg(s)
add_title(s, "AELSTM-Family: Substantial but Incomplete Recovery")
page_num(s, 3)
top, h = content_box()
add_picture_fit(s, f"{AELSTM}/leakage_diagnostic_2012/evergreen/leakage_r2_comparison.png",
                  MARGIN, top, Inches(7.7), h, align="left")
add_bullets(s, MARGIN + Inches(8.0), top + Inches(0.3), Inches(4.4), h, [
    "Mean RMSE ↓ ~46% across all 8 models once 2012 is included in fitting.",
    "Mean R² ↑ +5.17 — a consistent effect across very different architectures.",
    "RF recovers the most, reaching R² = +0.59 — the only model to cross into positive territory.",
    "The other 7 models improve substantially but remain negative (roughly −0.6 to −4.2).",
], size=15.5, space_after=14)
add_takeaway(s, "Every model improves — but only Random Forest crosses into positive R².")

# ============================================================ SLIDE 4: CHRONOS-2 RESULTS
s = add_slide(); set_bg(s)
add_title(s, "Chronos-2: LoRA Fine-Tuning With 2012 Included")
page_num(s, 4)
top, h = content_box()
add_picture_fit(s, f"{CHRONOS2}/leakage_diagnostic_2012/evergreen/leakage_prediction_curves_chronos2.png",
                  MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(3.5), align="center")
add_bullets(s, MARGIN, top + Inches(3.7), SLIDE_W - 2 * MARGIN, Inches(1.1), [
    "RMSE ↓ 59.9%   |   R²: −8.29 → −0.50   |   Pearson r: 0.13 → 0.65   |   ACC: −0.63 → 0.88",
    "Zero-shot has no trainable weights, so LoRA fine-tuning is the only way to let Chronos-2 "
    "“see” 2012; both conditions reuse the existing, already-validated LoRA hyperparameters.",
], size=16, space_after=10)
add_takeaway(s, "Chronos-2 shows the same pattern as AELSTM: a large recovery, but still short of a good fit.")

# ============================================================ SLIDE 5: ANSWER
s = add_slide(); set_bg(s)
add_rect(s, 0, 0, SLIDE_W, Inches(1.65), ACCENT)
add_text(s, MARGIN, Inches(0.3), Inches(12.2), Inches(0.4),
          [("ANSWER TO PROF. WANG'S QUESTION", 13, RGBColor(0xD8, 0xE6, 0xE1), True, False)])
add_text(s, MARGIN, Inches(0.62), Inches(12.2), Inches(0.95),
          [("Allowing the models to see 2012 substantially improves their ability "
            "to reproduce the 2012 LAI dynamics.", 22, WHITE, True, False)], line_spacing=1.1)
page_num(s, 5)

top = Inches(2.0)
items = [
    ("1", "The dominant cause of the original failure appears to be distribution shift.",
     "2012 was a real drought, unlike anything in the preceding training history."),
    ("2", "The effect is consistent across very different model architectures, including Chronos-2.",
     "RF, 6 neural sequence models, and a large pretrained transformer all improve substantially "
     "once 2012 participates in fitting — the signature of a shared external cause."),
    ("3", "Model-specific limitations still remain.",
     "Most models — including Chronos-2 — still have negative R² even after seeing 2012, so "
     "distribution shift is not the entire explanation."),
]
row_h = Inches(1.55)
for i, (num, headline, sub) in enumerate(items):
    y = top + i * row_h
    add_rect(s, MARGIN, y, Inches(0.55), Inches(0.55), ACCENT)
    add_text(s, MARGIN, y, Inches(0.55), Inches(0.55), [(num, 18, WHITE, True, False)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MARGIN + Inches(0.8), y - Inches(0.02), Inches(11.6), Inches(0.5),
              [(headline, 18, INK, True, False)], line_spacing=1.1)
    add_text(s, MARGIN + Inches(0.8), y + Inches(0.5), Inches(11.6), Inches(0.85),
              [(sub, 14.5, MUTED, False, False)], line_spacing=1.2)

out_path = f"{ROOT}/Chronos2-vegetation-forecasting/Leakage_Diagnostic_2012_Summary.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
